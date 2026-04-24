#!/usr/bin/env python3
"""
doc2pptx — Convert documents (txt, md, docx, pdf) into PowerPoint presentations.

Optionally accepts an existing .pptx as a style template to inherit backgrounds,
color themes, fonts, and slide layouts.

Usage:
    python doc2pptx.py input.md -o output.pptx
    python doc2pptx.py report.pdf -o output.pptx --template brand.pptx
    python doc2pptx.py notes.txt -o output.pptx --template brand.pptx --title "My Deck"
"""

from __future__ import annotations

import argparse
import copy
import logging
import os
import re
import sys
import time
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


logger = logging.getLogger("doc2pptx")


def _default_log_path(input_path: str | None) -> Path:
    """Build a timestamped default log path under ./logs/."""
    stem = Path(input_path).stem if input_path else "run"
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return Path("logs") / f"doc2pptx_{ts}_{stem}.log"


def configure_logging(
    log_file: str | Path | None = None,
    verbose: bool = False,
    quiet: bool = False,
) -> Path | None:
    """Configure terminal + optional file logging for the `doc2pptx` logger.

    Terminal: concise, INFO by default (WARNING if quiet, DEBUG if verbose).
    File (when `log_file` is set): detailed, DEBUG always — captures full
    document extraction previews and full LLM prompt/input/response exchanges.

    Returns the resolved log file path (or None if file logging disabled).
    """
    # Clear any handlers attached to our namespace from a previous run (Gradio
    # keeps the process alive and would otherwise duplicate output).
    for h in list(logger.handlers):
        logger.removeHandler(h)
    logger.setLevel(logging.DEBUG)
    logger.propagate = False

    term_level = logging.WARNING if quiet else (logging.DEBUG if verbose else logging.INFO)
    term = logging.StreamHandler(stream=sys.stderr)
    term.setLevel(term_level)
    term.setFormatter(logging.Formatter("%(message)s"))
    logger.addHandler(term)

    resolved: Path | None = None
    if log_file is not None:
        resolved = Path(log_file)
        resolved.parent.mkdir(parents=True, exist_ok=True)
        fh = logging.FileHandler(resolved, mode="w", encoding="utf-8")
        fh.setLevel(logging.DEBUG)
        fh.setFormatter(logging.Formatter(
            "%(asctime)s %(levelname)-5s %(message)s",
            datefmt="%H:%M:%S",
        ))
        logger.addHandler(fh)
        logger.debug("Log file opened at %s", resolved)
    return resolved


def _preview(text: str, n: int = 240) -> str:
    """Single-line preview of a longer string, for INFO-level logs."""
    text = (text or "").strip().replace("\n", " ⏎ ")
    if len(text) <= n:
        return text
    return text[:n] + f"… [+{len(text) - n} chars]"


# ─────────────────────────────────────────────────────────────
# 1. DATA MODEL — Intermediate representation of slide content
# ─────────────────────────────────────────────────────────────

@dataclass
class SlideContent:
    """One logical slide."""
    title: str = ""
    body_lines: list[str] = field(default_factory=list)
    level: int = 0  # heading depth (1 = H1/title, 2 = H2/section, etc.)
    slide_type: str = "content"  # "title", "section", "content", "table"
    # Table data (used when slide_type == "table")
    table_headers: list[str] = field(default_factory=list)
    table_rows: list[list[str]] = field(default_factory=list)


@dataclass
class DeckContent:
    """Full parsed document ready for slide generation."""
    title: str = ""
    subtitle: str = ""
    slides: list[SlideContent] = field(default_factory=list)


# ─────────────────────────────────────────────────────────────
# 2. DOCUMENT PARSERS — Extract structured content from files
# ─────────────────────────────────────────────────────────────

def read_txt(path: str) -> str:
    """Read plain text file."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def read_markdown(path: str) -> str:
    """Read markdown file (returns raw markdown text)."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def read_docx(path: str) -> str:
    """Extract text from a .docx file, preserving heading markers."""
    from docx import Document as DocxDocument
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

    doc = DocxDocument(path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        style_name = (para.style.name or "").lower()
        # Map Word heading styles to markdown-style markers
        if "heading 1" in style_name:
            lines.append(f"# {text}")
        elif "heading 2" in style_name:
            lines.append(f"## {text}")
        elif "heading 3" in style_name:
            lines.append(f"### {text}")
        elif "heading 4" in style_name:
            lines.append(f"#### {text}")
        elif "title" in style_name:
            lines.append(f"# {text}")
        elif "subtitle" in style_name:
            lines.append(f"_subtitle: {text}")
        elif "list" in style_name:
            lines.append(f"- {text}")
        else:
            lines.append(text)
    return "\n".join(lines)


def read_pdf(path: str) -> str:
    """Extract text from a PDF file."""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    texts.append(page_text)
        if texts:
            return "\n\n".join(texts)
    except Exception:
        pass

    # Fallback to pypdf
    from pypdf import PdfReader
    reader = PdfReader(path)
    texts = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            texts.append(page_text)
    return "\n\n".join(texts)


def read_html(path: str) -> str:
    """Extract text from an HTML file, mapping headings to markdown markers."""
    from bs4 import BeautifulSoup

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    lines = []
    for el in soup.find_all(["h1", "h2", "h3", "h4", "p", "li"]):
        text = el.get_text(strip=True)
        if not text:
            continue
        tag = el.name.lower()
        if tag == "h1":
            lines.append(f"# {text}")
        elif tag == "h2":
            lines.append(f"## {text}")
        elif tag == "h3":
            lines.append(f"### {text}")
        elif tag == "h4":
            lines.append(f"#### {text}")
        elif tag == "li":
            lines.append(f"- {text}")
        else:
            lines.append(text)
        lines.append("")  # blank line after each block element
    return "\n".join(lines)


# ─── Spreadsheet reader (returns DeckContent, not raw text) ──

MAX_TABLE_ROWS_PER_SLIDE = 15   # data rows per table slide (excl. header)
MAX_TABLE_COLS = 10             # columns before we truncate
MAX_COL_CHAR_WIDTH = 30        # truncate long cell values


def _sanitize_cell(val) -> str:
    """Convert a cell value to a clean display string."""
    if val is None:
        return ""
    import math
    if isinstance(val, float):
        if math.isnan(val):
            return ""
        # Show ints cleanly, floats with reasonable precision
        if val == int(val):
            return str(int(val))
        return f"{val:.2f}"
    s = str(val).strip()
    if len(s) > MAX_COL_CHAR_WIDTH:
        s = s[: MAX_COL_CHAR_WIDTH - 1] + "…"
    return s


def read_xlsx(path: str, deck_title: str = "", max_rows_per_slide: int = MAX_TABLE_ROWS_PER_SLIDE) -> DeckContent:
    """
    Read an xlsx/xls/csv/tsv file and produce a DeckContent with table slides.

    Each worksheet becomes one or more table slides. If a sheet has more rows
    than max_rows_per_slide, it is split across multiple slides.
    """
    import pandas as pd

    ext = Path(path).suffix.lower()

    # Read all sheets
    if ext in (".csv", ".tsv"):
        sep = "\t" if ext == ".tsv" else ","
        sheets = {"Sheet1": pd.read_csv(path, sep=sep)}
    else:
        sheets = pd.read_excel(path, sheet_name=None)

    file_stem = Path(path).stem
    deck = DeckContent(title=deck_title or file_stem)

    # Title slide
    deck.slides.append(SlideContent(
        title=deck.title,
        slide_type="title",
    ))

    for sheet_name, df in sheets.items():
        if df.empty:
            continue

        # Truncate columns if too many
        if len(df.columns) > MAX_TABLE_COLS:
            df = df.iloc[:, :MAX_TABLE_COLS]

        headers = [_sanitize_cell(c) for c in df.columns]

        # Convert all data to strings
        all_rows = []
        for _, row in df.iterrows():
            all_rows.append([_sanitize_cell(v) for v in row])

        # Chunk rows into slides
        for chunk_idx in range(0, max(1, len(all_rows)), max_rows_per_slide):
            chunk = all_rows[chunk_idx : chunk_idx + max_rows_per_slide]
            page_num = chunk_idx // max_rows_per_slide + 1
            total_pages = (len(all_rows) + max_rows_per_slide - 1) // max_rows_per_slide

            if total_pages == 1:
                slide_title = sheet_name
            else:
                slide_title = f"{sheet_name} ({page_num}/{total_pages})"

            deck.slides.append(SlideContent(
                title=slide_title,
                slide_type="table",
                table_headers=headers,
                table_rows=chunk,
            ))

    return deck


def read_document(path: str) -> str:
    """Route to the right reader based on file extension."""
    ext = Path(path).suffix.lower()
    readers = {
        ".txt": read_txt,
        ".md": read_markdown,
        ".markdown": read_markdown,
        ".docx": read_docx,
        ".pdf": read_pdf,
        ".html": read_html,
        ".htm": read_html,
    }
    reader = readers.get(ext)
    if reader is None:
        raise ValueError(
            f"Unsupported file type: '{ext}'. "
            f"Supported: {', '.join(readers.keys())}"
        )
    try:
        src_bytes = os.path.getsize(path)
    except OSError:
        src_bytes = -1
    t0 = time.monotonic()
    text = reader(path)
    dt = time.monotonic() - t0
    n_lines = text.count("\n") + (0 if text.endswith("\n") or not text else 1)
    logger.info(
        "   Extracted %d chars (%d lines) from %s file (%d bytes) in %.2fs",
        len(text), n_lines, ext, src_bytes, dt,
    )
    logger.info("   Preview: %s", _preview(text, 240))
    logger.debug("=== EXTRACTED TEXT (%d chars) ===\n%s\n=== END EXTRACTED TEXT ===", len(text), text)
    return text


# ─────────────────────────────────────────────────────────────
# 2b. LLM REWRITE (Ollama) — prose → PPT-friendly markdown
# ─────────────────────────────────────────────────────────────

DEFAULT_OLLAMA_HOST = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
DEFAULT_OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "gemma4:26b")

DEFAULT_REWRITE_PROMPT = """You are rewriting a document so it converts cleanly into a PowerPoint deck.

Rules:
- Output ONLY markdown. No preamble, no commentary, no code fences.
- Use `#` for the deck title, `##` for each new section title, `###` for slide headings. There should always be at least 2 slides per section, often more; group as many slides together under the same section as reasonably possible.
- Under each slide title, write 5-7 bullet points starting with `- `.
- Preserve the original meaning, facts, numbers, and ordering. Do not invent content.
- Do not leave any content out, but do summarize ideas instead of repeating verbatim.
- Drop filler, repetition, and boilerplate.
"""

CONTINUATION_REWRITE_PROMPT = """You are rewriting part of a longer document so it converts cleanly into a PowerPoint deck. This chunk is a CONTINUATION of an earlier rewrite.

Rules:
- Output ONLY markdown. No preamble, no commentary, no code fences.
- Do NOT emit a `#` deck title — the deck title was already produced. Start directly with a `##` slide title.
- Use `##` for each new major section title and `###` for slide headers. There should always be at least 2 slides per section, often more; group as many slides together under the same section as reasonably possible.
- Under each slide title, write 5-7 bullet points starting with `- `.
- Preserve the original meaning, facts, numbers, and ordering. Do not invent content.
- Do not leave any important content out, but do not repeat verbatim unless a direct quote is needed for clarity or completeness.
- Drop filler, repetition, boilerplate, etc.
"""


def rewrite_for_pptx(
    raw_text: str,
    host: str = DEFAULT_OLLAMA_HOST,
    model: str = DEFAULT_OLLAMA_MODEL,
    system_prompt: str | None = None,
    timeout: float = 240.0,
) -> str:
    """Ask a local Ollama server to reshape raw_text into PPT-friendly markdown.

    Returns the rewritten markdown. Raises on HTTP/connection errors — callers
    decide whether to fall back to the original text.
    """
    import httpx

    payload = {
        "model": model,
        "prompt": raw_text,
        "system": system_prompt or DEFAULT_REWRITE_PROMPT,
        "stream": False,
        "options": {"temperature": 0.2},
    }
    url = host.rstrip("/") + "/api/generate"
    logger.debug("POST %s (model=%s, prompt=%d chars, system=%d chars)",
                 url, model, len(raw_text), len(payload["system"]))
    with httpx.Client(timeout=timeout) as client:
        resp = client.post(url, json=payload)
        resp.raise_for_status()
        data = resp.json()
    return (data.get("response") or "").strip()


_HEADING_LINE_RE = re.compile(r"^#{1,4}\s+\S")


def _split_into_blocks(text: str) -> list[str]:
    """Split text into atomic blocks for greedy chunk packing.

    A block is either:
      - a heading line (starts with 1-4 `#` followed by whitespace and text), or
      - a paragraph (run of non-blank lines separated by blank lines).

    Blocks preserve their leading/trailing newlines so re-joining with `\\n\\n`
    produces readable text.
    """
    blocks: list[str] = []
    current: list[str] = []

    def flush() -> None:
        if current:
            joined = "\n".join(current).strip()
            if joined:
                blocks.append(joined)
            current.clear()

    for line in text.splitlines():
        if _HEADING_LINE_RE.match(line):
            # Headings are their own block.
            flush()
            blocks.append(line.strip())
        elif not line.strip():
            flush()
        else:
            current.append(line)
    flush()
    return blocks


def _split_oversized_block(block: str, max_chunk_chars: int) -> list[str]:
    """Split a single block that is larger than max_chunk_chars on sentence
    boundaries. Falls back to hard character slicing if sentences are still
    too long.
    """
    if len(block) <= max_chunk_chars:
        return [block]
    sentences = re.split(r"(?<=[.!?])\s+", block)
    pieces: list[str] = []
    buf = ""
    for s in sentences:
        if not s:
            continue
        if len(s) > max_chunk_chars:
            # Emit whatever is buffered, then hard-slice the long sentence.
            if buf:
                pieces.append(buf)
                buf = ""
            for i in range(0, len(s), max_chunk_chars):
                pieces.append(s[i : i + max_chunk_chars])
            continue
        candidate = (buf + " " + s).strip() if buf else s
        if len(candidate) > max_chunk_chars:
            pieces.append(buf)
            buf = s
        else:
            buf = candidate
    if buf:
        pieces.append(buf)
    return pieces


def _tail_sentences(text: str, n: int) -> str:
    """Return the last `n` sentences of text for use as cross-chunk overlap."""
    if n <= 0 or not text.strip():
        return ""
    sentences = re.split(r"(?<=[.!?])\s+", text.strip())
    sentences = [s for s in sentences if s]
    if not sentences:
        return ""
    return " ".join(sentences[-n:]).strip()


def chunk_document_for_llm(
    text: str,
    max_chunk_chars: int = 6000,
    overlap_sentences: int = 2,
) -> list[str]:
    """Split `text` into structure-aware chunks for LLM rewriting.

    Splitting priority: markdown headings (`#`..`####`) → blank-line paragraphs
    → sentences (only when a single paragraph exceeds `max_chunk_chars`).

    The last `overlap_sentences` sentences of chunk N are prepended to chunk
    N+1 as orientation context. The overlap is labelled so the model knows not
    to reproduce it. Always returns at least one chunk.
    """
    if max_chunk_chars <= 0:
        raise ValueError("max_chunk_chars must be positive")
    text = text or ""
    if not text.strip():
        return [text]

    blocks = _split_into_blocks(text)
    if not blocks:
        return [text]

    # Pre-split any block larger than the budget.
    expanded: list[str] = []
    for b in blocks:
        if len(b) > max_chunk_chars:
            expanded.extend(_split_oversized_block(b, max_chunk_chars))
        else:
            expanded.append(b)

    chunks: list[str] = []
    current_parts: list[str] = []
    current_len = 0
    for block in expanded:
        extra = len(block) + (2 if current_parts else 0)  # "\n\n" separator
        if current_parts and current_len + extra > max_chunk_chars:
            chunks.append("\n\n".join(current_parts))
            current_parts = [block]
            current_len = len(block)
        else:
            current_parts.append(block)
            current_len += extra
    if current_parts:
        chunks.append("\n\n".join(current_parts))

    if len(chunks) <= 1 or overlap_sentences <= 0:
        final_chunks = chunks
    else:
        # Prepend overlap to every chunk after the first.
        with_overlap: list[str] = [chunks[0]]
        for i in range(1, len(chunks)):
            tail = _tail_sentences(chunks[i - 1], overlap_sentences)
            if tail:
                header = f"[Context from previous section — do not repeat in output]\n{tail}"
                with_overlap.append(f"{header}\n\n{chunks[i]}")
            else:
                with_overlap.append(chunks[i])
        final_chunks = with_overlap

    # Log distribution so the user can decide whether to tune max_chunk_chars.
    sizes = [len(c) for c in final_chunks]
    if sizes:
        logger.info(
            "   Chunking: %d chunk(s) — size min=%d / avg=%d / max=%d (budget=%d, overlap=%d)",
            len(sizes), min(sizes), sum(sizes) // len(sizes), max(sizes),
            max_chunk_chars, overlap_sentences,
        )
        for i, c in enumerate(final_chunks, start=1):
            logger.debug(
                "--- CHUNK %d/%d (%d chars) ---\n%s\n--- END CHUNK %d ---",
                i, len(final_chunks), len(c), c, i,
            )
    return final_chunks


def _strip_deck_title_lines(markdown: str) -> str:
    """Remove any `# ...` deck-title lines. Used on continuation-chunk output
    as a safety net — the deck title must appear only once.
    """
    kept = [ln for ln in markdown.splitlines() if not re.match(r"^#\s+\S", ln)]
    return "\n".join(kept)


def _derive_continuation_prompt(base_prompt: str) -> str:
    """Given a (possibly user-customized) base prompt, produce a continuation
    variant that forbids emitting a new `#` deck title.
    """
    if base_prompt.strip() == DEFAULT_REWRITE_PROMPT.strip():
        return CONTINUATION_REWRITE_PROMPT
    addendum = (
        "\n\nIMPORTANT — this input is a CONTINUATION of an earlier rewrite: "
        "do NOT emit a `#` deck title. Start directly with a `##` slide title. "
        "A short excerpt from the previous section may appear at the top of the "
        "input as context — do not repeat it in the output."
    )
    return base_prompt.rstrip() + addendum


def _format_eta(seconds: float) -> str:
    seconds = max(0, int(seconds))
    if seconds < 60:
        return f"{seconds}s"
    if seconds < 3600:
        return f"{seconds // 60}m{seconds % 60:02d}s"
    h, rem = divmod(seconds, 3600)
    return f"{h}h{rem // 60:02d}m"


def rewrite_for_pptx_chunked(
    raw_text: str,
    host: str = DEFAULT_OLLAMA_HOST,
    model: str = DEFAULT_OLLAMA_MODEL,
    system_prompt: str | None = None,
    timeout: float = 240.0,
    max_chunk_chars: int = 6000,
    overlap_sentences: int = 2,
) -> str:
    """Chunk `raw_text`, rewrite each chunk via Ollama, and stitch results.

    The first chunk uses `system_prompt` (or `DEFAULT_REWRITE_PROMPT`).
    Subsequent chunks use a continuation variant that forbids a new `#` deck
    title. If any single chunk fails, its raw text is substituted so the
    pipeline degrades gracefully (mirroring the single-pass failure mode).
    """
    chunks = chunk_document_for_llm(
        raw_text,
        max_chunk_chars=max_chunk_chars,
        overlap_sentences=overlap_sentences,
    )

    base_prompt = system_prompt or DEFAULT_REWRITE_PROMPT
    continuation_prompt = _derive_continuation_prompt(base_prompt)
    total = len(chunks)

    logger.info("   Rewriting with %s at %s (%d chunk(s))", model, host, total)
    logger.debug("=== BASE SYSTEM PROMPT ===\n%s\n=== END BASE SYSTEM PROMPT ===", base_prompt)
    if total > 1:
        logger.debug(
            "=== CONTINUATION SYSTEM PROMPT ===\n%s\n=== END CONTINUATION SYSTEM PROMPT ===",
            continuation_prompt,
        )

    rewritten_parts: list[str] = []
    durations: list[float] = []
    wall_start = time.monotonic()
    for i, chunk in enumerate(chunks, start=1):
        role = "base" if i == 1 else "continuation"
        prompt = base_prompt if i == 1 else continuation_prompt
        logger.info("   [chunk %d/%d] sending %d chars (%s)…", i, total, len(chunk), role)
        logger.debug(
            "=== LLM REQUEST chunk %d/%d (%s) ===\n--- USER INPUT (%d chars) ---\n%s\n=== END REQUEST chunk %d ===",
            i, total, role, len(chunk), chunk, i,
        )

        t0 = time.monotonic()
        failed = False
        try:
            out = rewrite_for_pptx(
                chunk,
                host=host,
                model=model,
                system_prompt=prompt,
                timeout=timeout,
            )
        except Exception as exc:
            failed = True
            logger.warning(
                "   [chunk %d/%d] rewrite failed (%s); using raw chunk text.",
                i, total, exc,
            )
            logger.debug("=== LLM ERROR chunk %d/%d ===\n%r", i, total, exc)
            out = chunk
        dt = time.monotonic() - t0
        durations.append(dt)

        if not failed and not out.strip():
            logger.warning(
                "   [chunk %d/%d] model returned empty output; using raw chunk text.",
                i, total,
            )
            out = chunk

        if i > 1:
            before = out
            out = _strip_deck_title_lines(out)
            if out != before:
                logger.debug("   [chunk %d/%d] stripped stray `#` deck-title line(s) from output.", i, total)

        # Progress + ETA based on rolling average.
        avg = sum(durations) / len(durations)
        remaining = (total - i) * avg
        elapsed = time.monotonic() - wall_start
        logger.info(
            "   [chunk %d/%d] ← %d chars in %.1fs  (avg %.1fs/chunk, elapsed %s, ETA %s)",
            i, total, len(out), dt, avg,
            _format_eta(elapsed), _format_eta(remaining) if total > i else "done",
        )
        logger.debug(
            "=== LLM RESPONSE chunk %d/%d (%.2fs, %d chars) ===\n%s\n=== END RESPONSE chunk %d ===",
            i, total, dt, len(out), out, i,
        )

        rewritten_parts.append(out.strip())

    total_elapsed = time.monotonic() - wall_start
    logger.info("   Rewrite complete: %d chunk(s) in %s", total, _format_eta(total_elapsed))
    return "\n\n".join(p for p in rewritten_parts if p)


# ─────────────────────────────────────────────────────────────
# 3. TEXT → STRUCTURED SLIDES PARSER
# ─────────────────────────────────────────────────────────────

def _is_heading(line: str) -> tuple[int, str] | None:
    """Check if a line is a markdown-style heading. Returns (level, text) or None."""
    m = re.match(r"^(#{1,4})\s+(.+)$", line)
    if m:
        return len(m.group(1)), m.group(2).strip()
    return None


def _is_bullet(line: str) -> str | None:
    """Check if a line is a bullet point. Returns the text or None."""
    m = re.match(r"^\s*[-*•]\s+(.+)$", line)
    if m:
        return m.group(1).strip()
    # Numbered list
    m = re.match(r"^\s*\d+[.)]\s+(.+)$", line)
    if m:
        return m.group(1).strip()
    return None


def _looks_like_heading(line: str, prev_blank: bool, next_blank: bool) -> bool:
    # TODO: Make this function better
    """
    Heuristic: detect "implicit headings" in plain text.
    A short standalone line surrounded by blank lines that doesn't end with
    a period is likely a heading/section title.
    """
    stripped = line.strip()
    if not stripped:
        return False
    # Must be relatively short
    if len(stripped) > 80:
        return False
    # Must be surrounded by at least one blank line
    # TODO: Sometimes, extracting text from pdfs causes line spacing to be compressed, so this may not apply
    if not (prev_blank or next_blank):
        return False
    # Should NOT end with sentence-ending punctuation
    if stripped[-1] in ".!?;:,":
        return False
    # Should not be a bullet
    if _is_bullet(stripped):
        return False
    # Bonus: mostly title-case or all-caps is a strong signal
    words = stripped.split()
    if len(words) <= 8:
        return True
    return False


def _chunk_body_lines(lines: list[str], max_lines: int = 7) -> list[list[str]]:
    """Split a long list of body lines into chunks that fit on a slide."""
    if len(lines) <= max_lines:
        return [lines]
    chunks = []
    for i in range(0, len(lines), max_lines):
        chunks.append(lines[i : i + max_lines])
    return chunks


def _split_long_paragraph(text: str, max_chars: int = 200) -> list[str]:
    """Split a long paragraph into sentence-level bullet points."""
    sentences = re.split(r"(?<=[.!?])\s+", text)
    result = []
    current = ""
    for s in sentences:
        if current and len(current) + len(s) + 1 > max_chars:
            result.append(current.strip())
            current = s
        else:
            current = (current + " " + s).strip() if current else s
    if current:
        result.append(current.strip())
    return result


def _has_markdown_headings(text: str) -> bool:
    """Check whether the text contains any markdown-style headings."""
    return bool(re.search(r"^#{1,4}\s+.+", text, re.MULTILINE))


def parse_text_to_deck(raw_text: str, deck_title: str = "", max_bullets: int = 7) -> DeckContent:
    """
    Parse raw text (with optional markdown headings) into a DeckContent.

    Strategy:
    1. If markdown headings are present → use them as slide boundaries.
    2. Otherwise, detect implicit headings (short standalone lines) and
       paragraph blocks as sections.
    3. Bullet points and body text under a heading become that slide's body.
    4. Long slides are auto-chunked.
    """
    lines = raw_text.split("\n")
    deck = DeckContent(title=deck_title or "Presentation")

    # Check for subtitle marker (from docx parsing)
    for i, line in enumerate(lines):
        if line.startswith("_subtitle: "):
            deck.subtitle = line.replace("_subtitle: ", "").strip()
            lines[i] = ""
            break

    has_md = _has_markdown_headings(raw_text)

    if has_md:
        deck = _parse_markdown_structured(lines, deck, deck_title, max_bullets)
    else:
        deck = _parse_plain_text(lines, deck, deck_title, max_bullets)

    return deck


def _parse_markdown_structured(
    lines: list[str], deck: DeckContent, deck_title: str, max_bullets: int
) -> DeckContent:
    """Parse text that contains markdown headings."""
    # Find the deck title from the first H1 if not provided
    if not deck_title:
        for line in lines:
            heading = _is_heading(line)
            if heading and heading[0] == 1:
                deck.title = heading[1]
                break

    current_slide: SlideContent | None = None
    body_buffer: list[str] = []

    def flush_slide():
        nonlocal current_slide, body_buffer
        if current_slide is not None:
            current_slide.body_lines = body_buffer
            # Drop the empty H2 companion content slide when no bullets
            # followed the section divider (avoids duplicate titles).
            if (
                current_slide.level == 2
                and current_slide.slide_type == "content"
                and not body_buffer
            ):
                current_slide = None
                body_buffer = []
                return
            chunks = _chunk_body_lines(current_slide.body_lines, max_bullets)
            if len(chunks) == 1:
                deck.slides.append(current_slide)
            else:
                for idx, chunk in enumerate(chunks):
                    s = SlideContent(
                        title=current_slide.title + (f" (cont.)" if idx > 0 else ""),
                        body_lines=chunk,
                        level=current_slide.level,
                        slide_type=current_slide.slide_type,
                    )
                    deck.slides.append(s)
        current_slide = None
        body_buffer = []

    first_h1_seen = False
    for line in lines:
        stripped = line.strip()
        heading = _is_heading(stripped)

        if heading:
            level, text = heading
            flush_slide()
            if level == 1:
                if not first_h1_seen:
                    first_h1_seen = True
                    deck.title = deck.title or text
                    current_slide = SlideContent(title=text, level=1, slide_type="title")
                else:
                    current_slide = SlideContent(title=text, level=1, slide_type="section")
            elif level == 2:
                # H2 is always a section divider. Emit the section slide
                # immediately, then start a companion content slide with the
                # same title so any bullets before the next heading spill
                # onto it. If no bullets arrive, flush_slide drops the empty
                # companion to avoid a duplicate title slide.
                deck.slides.append(
                    SlideContent(title=text, level=2, slide_type="section")
                )
                current_slide = SlideContent(title=text, level=2, slide_type="content")
            else:
                current_slide = SlideContent(title=text, level=level, slide_type="content")
            continue

        bullet_text = _is_bullet(stripped)
        if bullet_text:
            if current_slide is None:
                current_slide = SlideContent(title="", slide_type="content")
            body_buffer.append(bullet_text)
            continue

        if stripped:
            if current_slide is None:
                current_slide = SlideContent(title="", slide_type="content")
            # If this is a long paragraph, split into sentences
            if len(stripped) > 200:
                body_buffer.extend(_split_long_paragraph(stripped))
            else:
                body_buffer.append(stripped)

    flush_slide()
    return deck


def _parse_plain_text(
    lines: list[str], deck: DeckContent, deck_title: str, max_bullets: int
) -> DeckContent:
    """
    Parse unstructured plain text by detecting implicit headings
    (short standalone lines between blank lines) and grouping
    body paragraphs under them.
    """
    # Step 1: Detect implicit headings using heuristics
    annotated: list[tuple[str, str]] = []  # (type, text): "heading", "bullet", "body", "blank"

    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            annotated.append(("blank", ""))
            continue

        bullet_text = _is_bullet(stripped)
        if bullet_text:
            annotated.append(("bullet", bullet_text))
            continue

        prev_blank = (i == 0) or (lines[i - 1].strip() == "")
        next_blank = (i == len(lines) - 1) or (lines[i + 1].strip() == "")

        if _looks_like_heading(stripped, prev_blank, next_blank):
            annotated.append(("heading", stripped))
        else:
            annotated.append(("body", stripped))

    # Step 2: Group into sections (heading → body/bullets until next heading)
    sections: list[tuple[str, list[str]]] = []
    current_heading = ""
    current_body: list[str] = []

    for kind, text in annotated:
        if kind == "heading":
            if current_heading or current_body:
                sections.append((current_heading, current_body))
            current_heading = text
            current_body = []
        elif kind == "bullet":
            current_body.append(text)
        elif kind == "body":
            # Split long paragraphs into digestible lines
            if len(text) > 200:
                current_body.extend(_split_long_paragraph(text))
            else:
                current_body.append(text)
        # blank → skip

    if current_heading or current_body:
        sections.append((current_heading, current_body))

    # Step 3: Build slides from sections
    if not sections:
        return deck

    # Use the first section heading (or first body text) as deck title
    first_heading = sections[0][0]
    first_body = sections[0][1]

    if not deck_title:
        if first_heading:
            deck.title = first_heading
        elif first_body:
            deck.title = first_body[0][:60]
        else:
            deck.title = "Presentation"

    # Title slide
    deck.slides.append(SlideContent(
        title=deck.title,
        slide_type="title",
    ))

    # If the first section was used as the title, still include its body as a content slide
    start = 0
    if first_heading:
        # First section heading was consumed as the deck title — skip its heading
        # but keep its body content if any
        if first_body:
            chunks = _chunk_body_lines(first_body, max_bullets)
            for idx, chunk in enumerate(chunks):
                deck.slides.append(SlideContent(
                    title=first_heading + (" (cont.)" if idx > 0 else ""),
                    body_lines=chunk,
                    slide_type="content",
                ))
        start = 1
    elif first_body:
        # No heading on first section — body was already noted, use it
        chunks = _chunk_body_lines(first_body, max_bullets)
        for idx, chunk in enumerate(chunks):
            deck.slides.append(SlideContent(
                title="" if idx == 0 else "(cont.)",
                body_lines=chunk,
                slide_type="content",
            ))
        start = 1

    for heading, body in sections[start:]:
        if not heading and not body:
            continue

        chunks = _chunk_body_lines(body, max_bullets) if body else [[]]
        for idx, chunk in enumerate(chunks):
            slide_type = "section" if not body and heading else "content"
            deck.slides.append(SlideContent(
                title=heading + (" (cont.)" if idx > 0 else ""),
                body_lines=chunk,
                level=2,
                slide_type=slide_type,
            ))

    return deck


# ─────────────────────────────────────────────────────────────
# 4. TEMPLATE ANALYZER — Extract style info from a .pptx
# ─────────────────────────────────────────────────────────────

@dataclass
class TemplateStyle:
    """Style information extracted from a template presentation."""
    presentation: Presentation
    # Layout indices by purpose
    title_layout_idx: int = 0
    section_layout_idx: int = 0
    content_layout_idx: int = 1
    # Whether a real Section Header layout was found (vs. a fallback).
    # When False, section slides use the default-style renderer so they
    # still look distinct from plain content slides.
    section_layout_found: bool = False
    # Fonts detected from template
    title_font: str = "Calibri"
    body_font: str = "Calibri"
    title_size: Pt = field(default_factory=lambda: Pt(36))
    section_title_size: Pt | None = None
    content_title_size: Pt | None = None
    body_size: Pt = field(default_factory=lambda: Pt(16))
    # Colors
    title_color: RGBColor | None = None
    body_color: RGBColor | None = None
    background_fill: any = None


def _ph_type_name(ph) -> str:
    """Short, readable name for a placeholder's type (for logging)."""
    try:
        t = ph.placeholder_format.type
    except Exception:
        return "?"
    if t is None:
        return "?"
    return getattr(t, "name", str(t))


def _fmt_size(s) -> str:
    """Format a Pt/EMU size as a readable 'Npt' string for logging."""
    if s is None:
        return "None"
    try:
        return f"{s.pt}pt"
    except AttributeError:
        return str(s)


def _describe_background(element, part) -> str:
    """Return a short description of the background of a layout or master.

    Inspects ``<p:cSld>/<p:bg>`` and reports whether the background is a
    solid fill, gradient, image (with the rId and resolved target part
    name), a theme bgRef, or inherited from a parent. Purely observational.
    """
    cSld = element.find(qn("p:cSld"))
    if cSld is None:
        return "no cSld"
    bg = cSld.find(qn("p:bg"))
    if bg is None:
        return "inherits from master/theme"
    bgPr = bg.find(qn("p:bgPr"))
    if bgPr is not None:
        blip = bgPr.find(qn("a:blipFill"))
        if blip is not None:
            inner = blip.find(qn("a:blip"))
            r_embed = inner.get(qn("r:embed")) if inner is not None else None
            target = None
            if r_embed is not None and part is not None:
                try:
                    target = str(part.related_parts[r_embed].partname)
                except Exception:
                    target = "<unresolved>"
            return f"image (rId={r_embed}, target={target})"
        if bgPr.find(qn("a:solidFill")) is not None:
            return "solid fill"
        if bgPr.find(qn("a:gradFill")) is not None:
            return "gradient fill"
        if bgPr.find(qn("a:pattFill")) is not None:
            return "pattern fill"
        return "bgPr (unknown fill)"
    bgRef = bg.find(qn("p:bgRef"))
    if bgRef is not None:
        return f"theme bgRef (idx={bgRef.get('idx')})"
    return "bg element (empty)"


def _log_layout_inventory(prs: Presentation) -> None:
    """Dump every slide layout and its placeholders at DEBUG level."""
    if not logger.isEnabledFor(logging.DEBUG):
        return
    logger.debug("   Layout inventory (%d layout(s)):", len(prs.slide_layouts))
    for idx, layout in enumerate(prs.slide_layouts):
        phs = list(layout.placeholders)
        logger.debug("     [%d] %r  (%d placeholder(s))", idx, layout.name, len(phs))
        for ph in phs:
            pf = ph.placeholder_format
            default_text = ""
            if ph.has_text_frame:
                default_text = (ph.text_frame.text or "").strip().replace("\n", " ")
                if len(default_text) > 60:
                    default_text = default_text[:57] + "..."
            logger.debug(
                "         idx=%s type=%s name=%r default_text=%r",
                pf.idx, _ph_type_name(ph), ph.name, default_text,
            )


def _log_background_info(prs: Presentation) -> None:
    """Dump background fill info for every master and layout at DEBUG level."""
    if not logger.isEnabledFor(logging.DEBUG):
        return
    logger.debug("   Background inventory:")
    for m_idx, master in enumerate(prs.slide_masters):
        desc = _describe_background(master.element, master.part)
        logger.debug("     master[%d] %r: %s", m_idx, master.name, desc)
        for l_idx, layout in enumerate(master.slide_layouts):
            desc = _describe_background(layout.element, layout.part)
            logger.debug("       layout[%d] %r: %s", l_idx, layout.name, desc)


def _find_best_layout(prs: Presentation, target: str) -> tuple[int, bool]:
    """Find the layout index that best matches a target purpose.

    Returns ``(index, matched)`` where ``matched`` is True when a layout
    whose name actually matched a keyword was found (vs. falling back to a
    generic index). Callers use ``matched`` to decide whether to trust the
    layout for decorative purposes — e.g. section slides fall back to the
    default renderer when no real Section Header layout exists.
    """
    layouts = prs.slide_layouts
    target_lower = target.lower()

    # Keywords to match layout names. "blank" is deliberately excluded from
    # "section" — a blank layout has no decoration and is indistinguishable
    # from a content slide, which defeats the purpose of a section divider.
    keyword_map = {
        "title": ["title slide", "title", "cover"],
        "section": ["section header", "section", "divider"],
        "content": [
            "title and content", "two content", "content",
            "title, content", "title and body", "text",
        ],
    }

    keywords = keyword_map.get(target_lower, [target_lower])
    for priority_kw in keywords:
        for idx, layout in enumerate(layouts):
            if layout.name and priority_kw in layout.name.lower():
                logger.debug(
                    "   layout match: target=%r -> idx=%d name=%r (keyword=%r)",
                    target, idx, layout.name, priority_kw,
                )
                return idx, True

    # Fallback: title=0, section=0, content=1 (or 0 if only one layout)
    fallback = {"title": 0, "section": 0, "content": min(1, len(layouts) - 1)}
    fb_idx = fallback.get(target_lower, 0)
    logger.debug(
        "   layout fallback: target=%r -> idx=%d (no keyword match; available=%s)",
        target, fb_idx, [l.name for l in layouts],
    )
    return fb_idx, False


def _layout_title_size(layout) -> Pt | None:
    """Return the size of the first sized run on the layout's title
    placeholder (idx=0), or None if the layout doesn't advertise one."""
    for ph in layout.placeholders:
        if ph.placeholder_format.idx != 0 or not ph.has_text_frame:
            continue
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    logger.debug(
                        "   title size from layout %r: %s", layout.name, _fmt_size(run.font.size),
                    )
                    return run.font.size
    logger.debug("   title size from layout %r: <none advertised>", layout.name)
    return None


def _extract_font_info(prs: Presentation) -> dict:
    """Scan existing slides/layouts to detect the fonts and colors in use."""
    info = {
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": Pt(36),
        "body_size": Pt(16),
        "title_color": None,
        "body_color": None,
    }

    # Scan slide layouts for placeholder fonts
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        if ph.placeholder_format.idx == 0:  # Title
                            if font.name:
                                if info["title_font"] != font.name:
                                    logger.debug(
                                        "   font detect (layout %r, title ph): title_font %r -> %r",
                                        layout.name, info["title_font"], font.name,
                                    )
                                info["title_font"] = font.name
                            if font.size:
                                if info["title_size"] != font.size:
                                    logger.debug(
                                        "   font detect (layout %r, title ph): title_size %s -> %s",
                                        layout.name, _fmt_size(info["title_size"]), _fmt_size(font.size),
                                    )
                                info["title_size"] = font.size
                            try:
                                if font.color and font.color.type is not None:
                                    if info["title_color"] != font.color.rgb:
                                        logger.debug(
                                            "   font detect (layout %r, title ph): title_color %s -> %s",
                                            layout.name, info["title_color"], font.color.rgb,
                                        )
                                    info["title_color"] = font.color.rgb
                            except (AttributeError, TypeError):
                                pass
                        elif ph.placeholder_format.idx == 1:  # Body
                            if font.name:
                                if info["body_font"] != font.name:
                                    logger.debug(
                                        "   font detect (layout %r, body ph): body_font %r -> %r",
                                        layout.name, info["body_font"], font.name,
                                    )
                                info["body_font"] = font.name
                            if font.size:
                                if info["body_size"] != font.size:
                                    logger.debug(
                                        "   font detect (layout %r, body ph): body_size %s -> %s",
                                        layout.name, _fmt_size(info["body_size"]), _fmt_size(font.size),
                                    )
                                info["body_size"] = font.size
                            try:
                                if font.color and font.color.type is not None:
                                    if info["body_color"] != font.color.rgb:
                                        logger.debug(
                                            "   font detect (layout %r, body ph): body_color %s -> %s",
                                            layout.name, info["body_color"], font.color.rgb,
                                        )
                                    info["body_color"] = font.color.rgb
                            except (AttributeError, TypeError):
                                pass

    # Also scan actual slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        if font.name and info["title_font"] == "Calibri":
                            logger.debug(
                                "   font detect (existing slide shape %r): title_font Calibri -> %r",
                                shape.name, font.name,
                            )
                            info["title_font"] = font.name
                        if font.size and info["title_size"] == Pt(36):
                            logger.debug(
                                "   font detect (existing slide shape %r): title_size 36pt -> %s",
                                shape.name, _fmt_size(font.size),
                            )
                            info["title_size"] = font.size

    return info


def analyze_template(template_path: str) -> TemplateStyle:
    """Load a .pptx template and extract its style information."""
    logger.info("🎨 Analyzing template: %s", template_path)
    prs = Presentation(template_path)
    logger.debug(
        "   Template has %d slide master(s), %d slide layout(s), %d existing slide(s)",
        len(prs.slide_masters), len(prs.slide_layouts), len(prs.slides),
    )
    _log_layout_inventory(prs)
    _log_background_info(prs)

    font_info = _extract_font_info(prs)

    title_idx, _ = _find_best_layout(prs, "title")
    section_idx, section_found = _find_best_layout(prs, "section")
    content_idx, _ = _find_best_layout(prs, "content")

    layouts = prs.slide_layouts
    section_title_size = (
        _layout_title_size(layouts[section_idx]) if section_found else None
    )
    content_title_size = _layout_title_size(layouts[content_idx])

    style = TemplateStyle(
        presentation=prs,
        title_layout_idx=title_idx,
        section_layout_idx=section_idx,
        content_layout_idx=content_idx,
        section_layout_found=section_found,
        title_font=font_info["title_font"],
        body_font=font_info["body_font"],
        title_size=font_info["title_size"],
        section_title_size=section_title_size,
        content_title_size=content_title_size,
        body_size=font_info["body_size"],
        title_color=font_info["title_color"],
        body_color=font_info["body_color"],
    )

    def _name(i: int) -> str:
        try:
            return layouts[i].name or "<unnamed>"
        except Exception:
            return "<oob>"

    logger.info(
        "   Template summary: title=idx %d %r, section=idx %d %r (%s), content=idx %d %r",
        title_idx, _name(title_idx),
        section_idx, _name(section_idx),
        "matched" if section_found else "fallback -> default renderer",
        content_idx, _name(content_idx),
    )
    logger.info(
        "   Template fonts: title=%r %s color=%s, body=%r %s color=%s",
        style.title_font, _fmt_size(style.title_size), style.title_color,
        style.body_font, _fmt_size(style.body_size), style.body_color,
    )
    if style.section_title_size is not None:
        logger.debug("   section title size override: %s", _fmt_size(style.section_title_size))
    if style.content_title_size is not None:
        logger.debug("   content title size override: %s", _fmt_size(style.content_title_size))

    return style


# ─────────────────────────────────────────────────────────────
# 5. SLIDE GENERATORS — Build the .pptx output
# ─────────────────────────────────────────────────────────────

# ----- Default (no template) style constants -----

DEFAULT_COLORS = {
    "bg_dark": RGBColor(0x1E, 0x27, 0x61),      # navy
    "bg_light": RGBColor(0xFF, 0xFF, 0xFF),       # white
    "title_on_dark": RGBColor(0xFF, 0xFF, 0xFF),  # white text
    "title_on_light": RGBColor(0x1E, 0x27, 0x61), # navy text
    "body": RGBColor(0x33, 0x33, 0x33),           # dark gray
    "accent": RGBColor(0xCA, 0xDC, 0xFC),         # ice blue
    "subtle": RGBColor(0x64, 0x74, 0x8B),         # muted gray
}
DEFAULT_TITLE_FONT = "Georgia"
DEFAULT_BODY_FONT = "Calibri"


def _apply_font(run, font_name: str, size: Pt, color: RGBColor | None = None, bold: bool = False):
    """Apply font properties to a text run."""
    run.font.name = font_name
    run.font.size = size
    if color:
        run.font.color.rgb = color
    run.font.bold = bold


def _add_title_slide_default(prs: Presentation, slide_content: SlideContent, subtitle: str = ""):
    """Create a title slide with the default built-in style."""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_dark"]

    # Title
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_content.title
        _apply_font(run, DEFAULT_TITLE_FONT, Pt(40), DEFAULT_COLORS["title_on_dark"], bold=True)

    # Subtitle
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = subtitle or ""
            _apply_font(run, DEFAULT_BODY_FONT, Pt(18), DEFAULT_COLORS["accent"])
            break


def _add_section_slide_default(prs: Presentation, slide_content: SlideContent):
    """Create a section divider slide with default style."""
    # Use a blank layout and build manually
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide_layout = prs.slide_layouts[layout_idx]
    logger.info(
        "   slide: type=section title=%r -> default section renderer (blank layout idx=%d %r, navy bg)",
        slide_content.title, layout_idx, slide_layout.name,
    )
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_dark"]

    # Accent line
    from pptx.util import Inches
    from pptx.shapes.autoshape import Shape
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(2.2), Inches(1.5), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DEFAULT_COLORS["accent"]
    shape.line.fill.background()

    # Section title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_content.title
    _apply_font(run, DEFAULT_TITLE_FONT, Pt(36), DEFAULT_COLORS["title_on_dark"], bold=True)


def _add_content_slide_default(prs: Presentation, slide_content: SlideContent):
    """Create a content slide with default style."""
    slide_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_light"]

    y_cursor = Inches(0.5)

    # Slide title
    if slide_content.title:
        txBox = slide.shapes.add_textbox(Inches(0.7), y_cursor, Inches(8.6), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_content.title
        _apply_font(run, DEFAULT_TITLE_FONT, Pt(28), DEFAULT_COLORS["title_on_light"], bold=True)
        y_cursor = Inches(1.4)

    # Body content
    if slide_content.body_lines:
        txBox = slide.shapes.add_textbox(
            Inches(0.7), y_cursor, Inches(8.6), Inches(4.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(slide_content.body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)

            # Add a bullet indicator
            p.level = 0
            pPr = p._pPr
            if pPr is None:
                from pptx.oxml.ns import qn
                pPr = p._p.get_or_add_pPr()

            run = p.add_run()
            run.text = f"  •  {line}"
            _apply_font(run, DEFAULT_BODY_FONT, Pt(15), DEFAULT_COLORS["body"])


# ─── Table slide helpers ──────────────────────────────────────

def _compute_col_widths(headers: list[str], rows: list[list[str]], total_width: float) -> list[float]:
    """
    Compute proportional column widths based on max content length per column.
    Returns widths in inches that sum to total_width.
    """
    n_cols = len(headers)
    if n_cols == 0:
        return []

    max_lens = []
    for c in range(n_cols):
        col_max = len(headers[c])
        for row in rows:
            if c < len(row):
                col_max = max(col_max, len(row[c]))
        max_lens.append(max(col_max, 3))  # minimum 3 chars

    total_chars = sum(max_lens)
    # Proportional, with a minimum width per column
    min_width = 0.8  # inches
    widths = []
    for ml in max_lens:
        w = max(min_width, (ml / total_chars) * total_width)
        widths.append(w)

    # Normalize so they sum to total_width
    scale = total_width / sum(widths)
    widths = [w * scale for w in widths]
    return widths


def _build_table_on_slide(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left: float,
    top: float,
    width: float,
    height: float,
    header_font: str = "Calibri",
    body_font: str = "Calibri",
    header_size: Pt = None,
    body_size: Pt = None,
    header_color: RGBColor | None = None,
    body_color: RGBColor | None = None,
    header_bg: RGBColor | None = None,
    stripe_bg: RGBColor | None = None,
):
    """
    Add a formatted table shape to a slide.
    Handles dynamic column widths and alternating row stripes.
    """
    if header_size is None:
        header_size = Pt(11)
    if body_size is None:
        body_size = Pt(10)
    if header_bg is None:
        header_bg = RGBColor(0x1E, 0x27, 0x61)  # navy
    if header_color is None:
        header_color = RGBColor(0xFF, 0xFF, 0xFF)  # white
    if body_color is None:
        body_color = RGBColor(0x33, 0x33, 0x33)
    if stripe_bg is None:
        stripe_bg = RGBColor(0xF0, 0xF4, 0xF8)  # very light blue-gray

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header row

    # Compute proportional column widths
    col_widths = _compute_col_widths(headers, rows, width)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    table = table_shape.table

    # Set column widths
    for c, w in enumerate(col_widths):
        table.columns[c].width = Inches(w)

    # Disable built-in banding (we'll do our own striping)
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is not None:
        tblPr.set("bandRow", "0")
        tblPr.set("bandCol", "0")
        tblPr.set("firstRow", "0")
        tblPr.set("lastRow", "0")

    # Helper to style a cell
    def _style_cell(cell, text, font_name, font_size, font_color, bg_color=None, bold=False):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        # Reduce internal margins for compact tables
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold

        if bg_color:
            from pptx.oxml.ns import qn
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn("a:solidFill"), {})
            srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": str(bg_color)})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

    # Header row
    for c, h in enumerate(headers):
        _style_cell(
            table.cell(0, c), h,
            header_font, header_size, header_color,
            bg_color=header_bg, bold=True,
        )

    # Data rows with alternating stripes
    for r, row_data in enumerate(rows):
        bg = stripe_bg if r % 2 == 1 else None
        for c in range(n_cols):
            val = row_data[c] if c < len(row_data) else ""
            _style_cell(
                table.cell(r + 1, c), val,
                body_font, body_size, body_color,
                bg_color=bg,
            )

    return table_shape


def _add_table_slide_default(prs: Presentation, slide_content: SlideContent):
    """Create a table slide with the default built-in style."""
    slide_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_light"]

    y_cursor = 0.4

    # Slide title
    if slide_content.title:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_cursor), Inches(9.0), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_content.title
        _apply_font(run, DEFAULT_TITLE_FONT, Pt(24), DEFAULT_COLORS["title_on_light"], bold=True)
        y_cursor = 1.1

    # Table
    n_rows = len(slide_content.table_rows)
    # Dynamically size: more rows → smaller row height
    available_height = 5.625 - y_cursor - 0.3  # slide height minus margins
    row_height_estimate = min(0.35, available_height / (n_rows + 1))
    table_height = row_height_estimate * (n_rows + 1)

    _build_table_on_slide(
        slide,
        headers=slide_content.table_headers,
        rows=slide_content.table_rows,
        left=0.5,
        top=y_cursor,
        width=9.0,
        height=table_height,
        header_font=DEFAULT_BODY_FONT,
        body_font=DEFAULT_BODY_FONT,
        header_bg=DEFAULT_COLORS["bg_dark"],
        header_color=DEFAULT_COLORS["title_on_dark"],
        body_color=DEFAULT_COLORS["body"],
        stripe_bg=RGBColor(0xEE, 0xF2, 0xF7),
    )


def _add_table_slide_from_template(prs: Presentation, style: TemplateStyle, slide_content: SlideContent):
    """Create a table slide using the template's layout and style."""
    layouts = prs.slide_layouts
    layout_idx = min(style.content_layout_idx, len(layouts) - 1)
    layout = layouts[layout_idx]
    n_rows = len(slide_content.table_rows)
    n_cols = len(slide_content.table_headers)
    logger.info(
        "   slide: type=table title=%r -> layout idx=%d %r (%d cols x %d rows)",
        slide_content.title, layout_idx, layout.name, n_cols, n_rows,
    )
    slide = prs.slides.add_slide(layout)

    # Set the title placeholder if available
    title_set = False
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = slide_content.title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = style.title_font
                    if style.title_color:
                        run.font.color.rgb = style.title_color
            title_set = True
            logger.debug(
                "     table title placeholder found: idx=0 name=%r (filled)", ph.name,
            )
            break
    if not title_set:
        logger.debug(
            "     no idx=0 title placeholder on table layout; table title not rendered as placeholder"
        )

    available_height = 4.0
    row_height_estimate = min(0.35, available_height / (n_rows + 1))
    table_height = row_height_estimate * (n_rows + 1)

    # Derive header background from template title color or fall back to navy
    header_bg = style.title_color or RGBColor(0x1E, 0x27, 0x61)

    _build_table_on_slide(
        slide,
        headers=slide_content.table_headers,
        rows=slide_content.table_rows,
        left=0.5,
        top=1.5,
        width=9.0,
        height=table_height,
        header_font=style.body_font,
        body_font=style.body_font,
        header_size=Pt(11),
        body_size=style.body_size if style.body_size else Pt(10),
        header_bg=header_bg,
        header_color=RGBColor(0xFF, 0xFF, 0xFF),
        body_color=style.body_color or RGBColor(0x33, 0x33, 0x33),
    )


def _set_placeholder_text(ph, text: str, fallback_font: str,
                          fallback_size: Pt | None, fallback_color: RGBColor | None):
    """Write ``text`` into a placeholder while preserving the layout's
    existing run formatting where possible.

    The layout placeholder already carries the template designer's font,
    size, weight, and color. Calling ``ph.text = ...`` blows all of that
    away. Instead, reuse the first paragraph's first run (which inherits
    layout styling via the run/paragraph/default chain) and only fill in
    ``font.name``/``font.size``/``font.color`` when the layout itself
    didn't provide one.
    """
    logger.debug(
        "     set placeholder text: idx=%d name=%r text=%r",
        ph.placeholder_format.idx, ph.name, text[:60],
    )
    tf = ph.text_frame
    # Keep the first paragraph (inherits layout pPr); drop the rest.
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    p = tf.paragraphs[0]
    # Preserve the first run if present (keeps its rPr); otherwise make one.
    if p.runs:
        run = p.runs[0]
        # Clear any additional runs so we end up with a single run.
        for extra in list(p.runs[1:]):
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
    run.text = text
    # Only fill in properties the layout didn't already supply.
    applied = []
    if not run.font.name and fallback_font:
        run.font.name = fallback_font
        applied.append(f"name={fallback_font!r}")
    if run.font.size is None and fallback_size:
        run.font.size = fallback_size
        applied.append(f"size={_fmt_size(fallback_size)}")
    try:
        has_color = run.font.color and run.font.color.type is not None
    except (AttributeError, TypeError):
        has_color = False
    if not has_color and fallback_color:
        run.font.color.rgb = fallback_color
        applied.append(f"color={fallback_color}")
    if applied:
        logger.debug(
            "     font fallback applied (layout did not supply): %s", ", ".join(applied),
        )


def _fill_body_placeholder(ph, lines: list[str], style: TemplateStyle):
    """Write ``lines`` into a body placeholder, reusing the layout's
    existing bullet/paragraph formatting for the first line so the
    template's indent and bullet glyph stay intact."""
    logger.debug(
        "     fill body placeholder: idx=%d name=%r %d line(s)",
        ph.placeholder_format.idx, ph.name, len(lines),
    )
    tf = ph.text_frame
    # Drop every paragraph past the first so we can rebuild from a clean slate.
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    first_p = tf.paragraphs[0]
    # Clear runs from the first paragraph but keep its pPr (bullet/indent).
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)

    fallback_counts = {"name": 0, "size": 0, "color": 0}
    for i, line in enumerate(lines):
        if i == 0:
            p = first_p
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if not run.font.name and style.body_font:
            run.font.name = style.body_font
            fallback_counts["name"] += 1
        if run.font.size is None and style.body_size:
            run.font.size = style.body_size
            fallback_counts["size"] += 1
        try:
            has_color = run.font.color and run.font.color.type is not None
        except (AttributeError, TypeError):
            has_color = False
        if not has_color and style.body_color:
            run.font.color.rgb = style.body_color
            fallback_counts["color"] += 1
    if any(fallback_counts.values()):
        logger.debug(
            "     body font fallback applied on %d/%d lines: name=%d size=%d color=%d",
            max(fallback_counts.values()), len(lines),
            fallback_counts["name"], fallback_counts["size"], fallback_counts["color"],
        )


def _add_slide_from_template(prs: Presentation, style: TemplateStyle, slide_content: SlideContent, subtitle: str = ""):
    """Create a slide using the template's layouts and styles."""
    # Section slides without a real Section Header layout would be
    # indistinguishable from content slides; fall back to the default
    # renderer so they still look like section dividers.
    if slide_content.slide_type == "section" and not style.section_layout_found:
        logger.debug(
            "   section slide %r -> default renderer (no Section Header layout in template)",
            slide_content.title,
        )
        _add_section_slide_default(prs, slide_content)
        return

    layouts = prs.slide_layouts

    if slide_content.slide_type == "title":
        layout_idx = style.title_layout_idx
        title_size = style.title_size
    elif slide_content.slide_type == "section":
        layout_idx = style.section_layout_idx
        title_size = style.section_title_size or style.title_size
    else:
        layout_idx = style.content_layout_idx
        title_size = style.content_title_size or style.title_size

    layout_idx = min(layout_idx, len(layouts) - 1)
    layout = layouts[layout_idx]
    logger.info(
        "   slide: type=%s title=%r -> layout idx=%d %r (title_size=%s)",
        slide_content.slide_type, slide_content.title, layout_idx, layout.name, _fmt_size(title_size),
    )
    slide = prs.slides.add_slide(layout)

    # Track what we did with each placeholder so "left-untouched" cases
    # (i.e. stale empty layout boxes) show up in the log.
    ph_actions: list[str] = []
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title placeholder
            _set_placeholder_text(
                ph, slide_content.title,
                fallback_font=style.title_font,
                fallback_size=title_size,
                fallback_color=style.title_color,
            )
            ph_actions.append(
                f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=filled-title"
            )
        elif idx == 1:  # Body/subtitle placeholder
            if slide_content.slide_type == "title" and subtitle:
                _set_placeholder_text(
                    ph, subtitle,
                    fallback_font=style.body_font,
                    fallback_size=style.body_size,
                    fallback_color=style.body_color,
                )
                ph_actions.append(
                    f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=filled-subtitle"
                )
            elif slide_content.body_lines:
                _fill_body_placeholder(ph, slide_content.body_lines, style)
                ph_actions.append(
                    f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=filled-body"
                )
            else:
                ph_actions.append(
                    f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=LEFT-UNTOUCHED (no body/subtitle content)"
                )
        else:
            ph_actions.append(
                f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=LEFT-UNTOUCHED (no handler for this idx)"
            )

    if logger.isEnabledFor(logging.DEBUG):
        if ph_actions:
            logger.debug("     placeholder actions:")
            for line in ph_actions:
                logger.debug("       %s", line)
        else:
            logger.debug("     placeholder actions: (layout has no placeholders)")

    # Fallbacks for minimal layouts that have no title/body placeholders.
    # Without these, H3 slide titles vanish silently because there's nowhere
    # for the title to land.
    title_ph_found = any(ph.placeholder_format.idx == 0 for ph in slide.placeholders)
    body_ph_found = any(ph.placeholder_format.idx == 1 for ph in slide.placeholders)

    body_top = Inches(1.8)
    if not title_ph_found and slide_content.title and slide_content.slide_type != "section":
        logger.debug(
            "     fallback: layout has no title placeholder; adding custom title textbox"
        )
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(1.0))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        trun = tp.add_run()
        trun.text = slide_content.title
        trun.font.name = style.title_font
        trun.font.size = title_size or Pt(28)
        trun.font.bold = True
        if style.title_color:
            trun.font.color.rgb = style.title_color
        body_top = Inches(1.6)

    if not body_ph_found and slide_content.body_lines and slide_content.slide_type == "content":
        logger.debug(
            "     fallback: layout has no body placeholder; adding custom body textbox (%d line(s))",
            len(slide_content.body_lines),
        )
        txBox = slide.shapes.add_textbox(Inches(0.7), body_top, Inches(8.6), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(slide_content.body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"•  {line}"
            run.font.name = style.body_font
            run.font.size = style.body_size
            if style.body_color:
                run.font.color.rgb = style.body_color
            p.space_after = Pt(6)


# ─────────────────────────────────────────────────────────────
# 6. MAIN ORCHESTRATOR
# ─────────────────────────────────────────────────────────────

def generate_pptx(
    input_path: str,
    output_path: str,
    template_path: str | None = None,
    title: str | None = None,
    max_bullets: int = 7,
    max_table_rows: int = MAX_TABLE_ROWS_PER_SLIDE,
    use_llm: bool = True,
    ollama_host: str = DEFAULT_OLLAMA_HOST,
    ollama_model: str = DEFAULT_OLLAMA_MODEL,
    llm_prompt_file: str | None = None,
    max_chunk_chars: int = 20000,
    chunk_overlap_sentences: int = 2,
):
    """
    Main entry point: read a document file and produce a .pptx presentation.

    Args:
        input_path:    Path to the input document (.txt, .md, .docx, .pdf, .html, .xlsx, .csv)
        output_path:   Path for the generated .pptx file
        template_path: Optional path to a .pptx template for styling
        title:         Optional override for the presentation title
        max_bullets:   Maximum bullet points per slide before auto-splitting
        max_table_rows: Maximum data rows per table slide (xlsx/csv only)
        use_llm:       If True, rewrite extracted text with a local Ollama server
                       before parsing. Skipped automatically for spreadsheets and
                       when the server is unreachable.
        ollama_host:   Base URL of the Ollama server.
        ollama_model:  Model name to request from Ollama.
        llm_prompt_file: Optional path to a text file containing a custom system
                       prompt to override the default rewrite instructions.
        max_chunk_chars: Maximum characters per LLM chunk. Large inputs are
                       always split on headings/paragraphs so small local models
                       don't overflow their context window.
        chunk_overlap_sentences: Number of trailing sentences from the previous
                       chunk to prepend to the next as orientation context.
    """
    ext = Path(input_path).suffix.lower()
    is_spreadsheet = ext in (".xlsx", ".xls", ".xlsm", ".csv", ".tsv")

    # 1. Read & parse
    logger.info("📄 Reading %s…", input_path)

    if is_spreadsheet:
        deck = read_xlsx(input_path, deck_title=title or "", max_rows_per_slide=max_table_rows)
    else:
        raw_text = read_document(input_path)
        if not raw_text.strip():
            raise ValueError(f"No text content could be extracted from '{input_path}'.")

        if use_llm:
            logger.info("🧠 Rewriting with Ollama (%s) at %s…", ollama_model, ollama_host)
            try:
                system_prompt = None
                if llm_prompt_file:
                    system_prompt = Path(llm_prompt_file).read_text(encoding="utf-8")
                    logger.info("   Using custom system prompt from %s (%d chars)",
                                llm_prompt_file, len(system_prompt))
                rewritten = rewrite_for_pptx_chunked(
                    raw_text,
                    host=ollama_host,
                    model=ollama_model,
                    system_prompt=system_prompt,
                    max_chunk_chars=max_chunk_chars,
                    overlap_sentences=chunk_overlap_sentences,
                )
                if rewritten:
                    raw_text = rewritten
                    logger.info("   LLM rewrite applied (%d chars).", len(rewritten))
                    logger.debug("=== FINAL STITCHED MARKDOWN (%d chars) ===\n%s\n=== END STITCHED MARKDOWN ===",
                                 len(rewritten), rewritten)
                else:
                    logger.warning("   ⚠️  LLM returned empty output; using original text.")
            except Exception as exc:
                logger.warning("   ⚠️  LLM rewrite skipped (%s); using original text.", exc)

        logger.info("🔍 Parsing document structure…")
        deck = parse_text_to_deck(raw_text, deck_title=title or "", max_bullets=max_bullets)

    if not deck.slides:
        raise ValueError("Could not extract any meaningful content for slides.")

    table_count = sum(1 for s in deck.slides if s.slide_type == "table")
    other_count = len(deck.slides) - table_count
    parts = []
    if other_count:
        parts.append(f"{other_count} content")
    if table_count:
        parts.append(f"{table_count} table")
    logger.info("   Found %d slides (%s)", len(deck.slides), ", ".join(parts))
    for idx, sc in enumerate(deck.slides, start=1):
        logger.debug("   slide %d: type=%s level=%d title=%r bullets=%d",
                     idx, sc.slide_type, sc.level, sc.title, len(sc.body_lines))

    # 2. Build the presentation
    if template_path:
        logger.info("🎨 Loading template from %s…", template_path)
        style = analyze_template(template_path)
        prs = Presentation(template_path)

        # Remove any existing slides from the template
        for _ in range(len(prs.slides)):
            sldId = prs.slides._sldIdLst[0]
            rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId is None:
                for attr_key, attr_val in sldId.attrib.items():
                    if attr_key.endswith('}id') or attr_key == 'r:id':
                        rId = attr_val
                        break
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except KeyError:
                    pass
            prs.slides._sldIdLst.remove(sldId)

        logger.info("   Generating slides with template styling…")
        for sc in deck.slides:
            if sc.slide_type == "table":
                _add_table_slide_from_template(prs, style, sc)
            else:
                _add_slide_from_template(prs, style, sc, subtitle=deck.subtitle)
    else:
        logger.info("🎨 Using default styling…")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        for sc in deck.slides:
            if sc.slide_type == "title":
                _add_title_slide_default(prs, sc, subtitle=deck.subtitle)
            elif sc.slide_type == "section":
                _add_section_slide_default(prs, sc)
            elif sc.slide_type == "table":
                _add_table_slide_default(prs, sc)
            else:
                _add_content_slide_default(prs, sc)

    # 3. Save
    prs.save(output_path)
    logger.info("✅ Saved to %s (%d slides)", output_path, len(deck.slides))
    return output_path


# ─────────────────────────────────────────────────────────────
# 7. CLI
# ─────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Convert documents (txt, md, docx, pdf, html, xlsx, csv) to PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python doc2pptx.py notes.md -o presentation.pptx
  python doc2pptx.py report.pdf -o deck.pptx --template brand.pptx
  python doc2pptx.py article.docx -o slides.pptx --title "Q4 Report"
  python doc2pptx.py page.html -o slides.pptx --max-bullets 5
  python doc2pptx.py data.xlsx -o tables.pptx
  python doc2pptx.py data.csv -o tables.pptx --max-table-rows 10
  python doc2pptx.py report.pdf -o deck.pptx --no-llm
  python doc2pptx.py report.pdf -o deck.pptx --ollama-model qwen2.5
        """,
    )
    parser.add_argument("input", help="Input document path (.txt, .md, .docx, .pdf, .html, .xlsx, .csv, .tsv)")
    parser.add_argument("-o", "--output", default="output.pptx", help="Output .pptx path (default: output.pptx)")
    parser.add_argument("-t", "--template", default=None, help="Template .pptx to use for styling")
    parser.add_argument("--title", default=None, help="Override the presentation title")
    parser.add_argument("--max-bullets", type=int, default=7, help="Max bullet points per slide before splitting (default: 7)")
    parser.add_argument("--max-table-rows", type=int, default=MAX_TABLE_ROWS_PER_SLIDE,
                        help=f"Max data rows per table slide (default: {MAX_TABLE_ROWS_PER_SLIDE})")
    parser.add_argument("--no-llm", action="store_true",
                        help="Skip the Ollama rewrite step (LLM rewrite is on by default).")
    parser.add_argument("--ollama-host", default=DEFAULT_OLLAMA_HOST,
                        help=f"Ollama server base URL (default: {DEFAULT_OLLAMA_HOST}, env: OLLAMA_HOST)")
    parser.add_argument("--ollama-model", default=DEFAULT_OLLAMA_MODEL,
                        help=f"Ollama model name (default: {DEFAULT_OLLAMA_MODEL}, env: OLLAMA_MODEL)")
    parser.add_argument("--prompt-file", default=None,
                        help="Path to a text file containing a custom system prompt for the rewrite step.")
    parser.add_argument("--max-chunk-chars", type=int, default=6000,
                        help="Max characters per LLM chunk; large docs are split on headings/paragraphs (default: 6000).")
    parser.add_argument("--chunk-overlap-sentences", type=int, default=2,
                        help="Trailing sentences from the previous chunk to prepend to the next as context (default: 2).")
    parser.add_argument("--log-file", default=None,
                        help="Path for the verbose run log. If omitted, a timestamped log is written under ./logs/. "
                             "Full LLM prompts, inputs, and responses are captured at DEBUG level.")
    parser.add_argument("--no-log-file", action="store_true",
                        help="Disable writing a run log file (terminal logging still runs).")
    parser.add_argument("-v", "--verbose", action="store_true",
                        help="Show DEBUG-level output in the terminal (chunk previews, prompts, responses).")
    parser.add_argument("-q", "--quiet", action="store_true",
                        help="Only show WARNING-level output in the terminal (the log file still captures everything).")

    args = parser.parse_args()

    if args.verbose and args.quiet:
        print("Error: --verbose and --quiet are mutually exclusive.", file=sys.stderr)
        sys.exit(2)

    log_file_path: Path | None = None
    if not args.no_log_file:
        log_file_path = Path(args.log_file) if args.log_file else _default_log_path(args.input)
    resolved_log = configure_logging(
        log_file=log_file_path,
        verbose=args.verbose,
        quiet=args.quiet,
    )
    if resolved_log is not None:
        logger.info("📝 Writing run log to %s", resolved_log)

    if not os.path.exists(args.input):
        logger.error("Input file not found: %s", args.input)
        sys.exit(1)

    if args.template and not os.path.exists(args.template):
        logger.error("Template file not found: %s", args.template)
        sys.exit(1)

    if args.prompt_file and not os.path.exists(args.prompt_file):
        logger.error("Prompt file not found: %s", args.prompt_file)
        sys.exit(1)

    generate_pptx(
        input_path=args.input,
        output_path=args.output,
        template_path=args.template,
        title=args.title,
        max_bullets=args.max_bullets,
        max_table_rows=args.max_table_rows,
        use_llm=not args.no_llm,
        ollama_host=args.ollama_host,
        ollama_model=args.ollama_model,
        llm_prompt_file=args.prompt_file,
        max_chunk_chars=args.max_chunk_chars,
        chunk_overlap_sentences=args.chunk_overlap_sentences,
    )


if __name__ == "__main__":
    main()
