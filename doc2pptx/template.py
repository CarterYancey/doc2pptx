"""Template analyzer — extract style info from a .pptx template."""

from __future__ import annotations

import logging

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

from .logging_config import logger
from .models import TemplateStyle


def _ph_type_name(ph) -> str:
    """Short, readable name for a placeholder's type (for logging)."""
    try:
        t = ph.placeholder_format.type
    except Exception:
        return "?"
    if t is None:
        return "?"
    return getattr(t, "name", str(t))


def _fmt_size(s) -> str:
    """Format a Pt/EMU size as a readable 'Npt' string for logging."""
    if s is None:
        return "None"
    try:
        return f"{s.pt}pt"
    except AttributeError:
        return str(s)


def _describe_background(element, part) -> str:
    """Return a short description of the background of a layout or master.

    Inspects ``<p:cSld>/<p:bg>`` and reports whether the background is a
    solid fill, gradient, image (with the rId and resolved target part
    name), a theme bgRef, or inherited from a parent. Purely observational.
    """
    cSld = element.find(qn("p:cSld"))
    if cSld is None:
        return "no cSld"
    bg = cSld.find(qn("p:bg"))
    if bg is None:
        return "inherits from master/theme"
    bgPr = bg.find(qn("p:bgPr"))
    if bgPr is not None:
        blip = bgPr.find(qn("a:blipFill"))
        if blip is not None:
            inner = blip.find(qn("a:blip"))
            r_embed = inner.get(qn("r:embed")) if inner is not None else None
            target = None
            if r_embed is not None and part is not None:
                try:
                    target = str(part.related_parts[r_embed].partname)
                except Exception:
                    target = "<unresolved>"
            return f"image (rId={r_embed}, target={target})"
        if bgPr.find(qn("a:solidFill")) is not None:
            return "solid fill"
        if bgPr.find(qn("a:gradFill")) is not None:
            return "gradient fill"
        if bgPr.find(qn("a:pattFill")) is not None:
            return "pattern fill"
        return "bgPr (unknown fill)"
    bgRef = bg.find(qn("p:bgRef"))
    if bgRef is not None:
        return f"theme bgRef (idx={bgRef.get('idx')})"
    return "bg element (empty)"


def _log_layout_inventory(prs: Presentation) -> None:
    """Dump every slide layout and its placeholders at DEBUG level."""
    if not logger.isEnabledFor(logging.DEBUG):
        return
    logger.debug("   Layout inventory (%d layout(s)):", len(prs.slide_layouts))
    for idx, layout in enumerate(prs.slide_layouts):
        phs = list(layout.placeholders)
        logger.debug("     [%d] %r  (%d placeholder(s))", idx, layout.name, len(phs))
        for ph in phs:
            pf = ph.placeholder_format
            default_text = ""
            if ph.has_text_frame:
                default_text = (ph.text_frame.text or "").strip().replace("\n", " ")
                if len(default_text) > 60:
                    default_text = default_text[:57] + "..."
            logger.debug(
                "         idx=%s type=%s name=%r default_text=%r",
                pf.idx, _ph_type_name(ph), ph.name, default_text,
            )


def _log_background_info(prs: Presentation) -> None:
    """Dump background fill info for every master and layout at DEBUG level."""
    if not logger.isEnabledFor(logging.DEBUG):
        return
    logger.debug("   Background inventory:")
    for m_idx, master in enumerate(prs.slide_masters):
        desc = _describe_background(master.element, master.part)
        logger.debug("     master[%d] %r: %s", m_idx, master.name, desc)
        for l_idx, layout in enumerate(master.slide_layouts):
            desc = _describe_background(layout.element, layout.part)
            logger.debug("       layout[%d] %r: %s", l_idx, layout.name, desc)


def _find_best_layout(prs: Presentation, target: str) -> tuple[int, bool]:
    """Find the layout index that best matches a target purpose.

    Returns ``(index, matched)`` where ``matched`` is True when a layout
    whose name actually matched a keyword was found (vs. falling back to a
    generic index). Callers use ``matched`` to decide whether to trust the
    layout for decorative purposes — e.g. section slides fall back to the
    default renderer when no real Section Header layout exists.
    """
    layouts = prs.slide_layouts
    target_lower = target.lower()

    # Keywords to match layout names. "blank" is deliberately excluded from
    # "section" — a blank layout has no decoration and is indistinguishable
    # from a content slide, which defeats the purpose of a section divider.
    keyword_map = {
        "title": ["title slide", "title", "cover"],
        "section": ["section header", "section", "divider"],
        "content": [
            "title and content", "two content", "content",
            "title, content", "title and body", "text",
        ],
    }

    keywords = keyword_map.get(target_lower, [target_lower])
    for priority_kw in keywords:
        for idx, layout in enumerate(layouts):
            if layout.name and priority_kw in layout.name.lower():
                logger.debug(
                    "   layout match: target=%r -> idx=%d name=%r (keyword=%r)",
                    target, idx, layout.name, priority_kw,
                )
                return idx, True

    # Fallback: title=0, section=0, content=1 (or 0 if only one layout)
    fallback = {"title": 0, "section": 0, "content": min(1, len(layouts) - 1)}
    fb_idx = fallback.get(target_lower, 0)
    logger.debug(
        "   layout fallback: target=%r -> idx=%d (no keyword match; available=%s)",
        target, fb_idx, [l.name for l in layouts],
    )
    return fb_idx, False


def _layout_title_size(layout) -> Pt | None:
    """Return the size of the first sized run on the layout's title
    placeholder (idx=0), or None if the layout doesn't advertise one."""
    for ph in layout.placeholders:
        if ph.placeholder_format.idx != 0 or not ph.has_text_frame:
            continue
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    logger.debug(
                        "   title size from layout %r: %s", layout.name, _fmt_size(run.font.size),
                    )
                    return run.font.size
    logger.debug("   title size from layout %r: <none advertised>", layout.name)
    return None


def _extract_font_info(prs: Presentation) -> dict:
    """Scan existing slides/layouts to detect the fonts and colors in use."""
    info = {
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": Pt(36),
        "body_size": Pt(16),
        "title_color": None,
        "body_color": None,
    }

    # Scan slide layouts for placeholder fonts
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        if ph.placeholder_format.idx == 0:  # Title
                            if font.name:
                                if info["title_font"] != font.name:
                                    logger.debug(
                                        "   font detect (layout %r, title ph): title_font %r -> %r",
                                        layout.name, info["title_font"], font.name,
                                    )
                                info["title_font"] = font.name
                            if font.size:
                                if info["title_size"] != font.size:
                                    logger.debug(
                                        "   font detect (layout %r, title ph): title_size %s -> %s",
                                        layout.name, _fmt_size(info["title_size"]), _fmt_size(font.size),
                                    )
                                info["title_size"] = font.size
                            try:
                                if font.color and font.color.type is not None:
                                    if info["title_color"] != font.color.rgb:
                                        logger.debug(
                                            "   font detect (layout %r, title ph): title_color %s -> %s",
                                            layout.name, info["title_color"], font.color.rgb,
                                        )
                                    info["title_color"] = font.color.rgb
                            except (AttributeError, TypeError):
                                pass
                        elif ph.placeholder_format.idx == 1:  # Body
                            if font.name:
                                if info["body_font"] != font.name:
                                    logger.debug(
                                        "   font detect (layout %r, body ph): body_font %r -> %r",
                                        layout.name, info["body_font"], font.name,
                                    )
                                info["body_font"] = font.name
                            if font.size:
                                if info["body_size"] != font.size:
                                    logger.debug(
                                        "   font detect (layout %r, body ph): body_size %s -> %s",
                                        layout.name, _fmt_size(info["body_size"]), _fmt_size(font.size),
                                    )
                                info["body_size"] = font.size
                            try:
                                if font.color and font.color.type is not None:
                                    if info["body_color"] != font.color.rgb:
                                        logger.debug(
                                            "   font detect (layout %r, body ph): body_color %s -> %s",
                                            layout.name, info["body_color"], font.color.rgb,
                                        )
                                    info["body_color"] = font.color.rgb
                            except (AttributeError, TypeError):
                                pass

    # Also scan actual slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        if font.name and info["title_font"] == "Calibri":
                            logger.debug(
                                "   font detect (existing slide shape %r): title_font Calibri -> %r",
                                shape.name, font.name,
                            )
                            info["title_font"] = font.name
                        if font.size and info["title_size"] == Pt(36):
                            logger.debug(
                                "   font detect (existing slide shape %r): title_size 36pt -> %s",
                                shape.name, _fmt_size(font.size),
                            )
                            info["title_size"] = font.size

    return info


def analyze_template(template_path: str) -> TemplateStyle:
    """Load a .pptx template and extract its style information."""
    logger.info("🎨 Analyzing template: %s", template_path)
    prs = Presentation(template_path)
    logger.debug(
        "   Template has %d slide master(s), %d slide layout(s), %d existing slide(s)",
        len(prs.slide_masters), len(prs.slide_layouts), len(prs.slides),
    )
    _log_layout_inventory(prs)
    _log_background_info(prs)

    font_info = _extract_font_info(prs)

    title_idx, _ = _find_best_layout(prs, "title")
    section_idx, section_found = _find_best_layout(prs, "section")
    content_idx, _ = _find_best_layout(prs, "content")

    layouts = prs.slide_layouts
    section_title_size = (
        _layout_title_size(layouts[section_idx]) if section_found else None
    )
    content_title_size = _layout_title_size(layouts[content_idx])

    style = TemplateStyle(
        presentation=prs,
        title_layout_idx=title_idx,
        section_layout_idx=section_idx,
        content_layout_idx=content_idx,
        section_layout_found=section_found,
        title_font=font_info["title_font"],
        body_font=font_info["body_font"],
        title_size=font_info["title_size"],
        section_title_size=section_title_size,
        content_title_size=content_title_size,
        body_size=font_info["body_size"],
        title_color=font_info["title_color"],
        body_color=font_info["body_color"],
    )

    def _name(i: int) -> str:
        try:
            return layouts[i].name or "<unnamed>"
        except Exception:
            return "<oob>"

    logger.info(
        "   Template summary: title=idx %d %r, section=idx %d %r (%s), content=idx %d %r",
        title_idx, _name(title_idx),
        section_idx, _name(section_idx),
        "matched" if section_found else "fallback -> default renderer",
        content_idx, _name(content_idx),
    )
    logger.info(
        "   Template fonts: title=%r %s color=%s, body=%r %s color=%s",
        style.title_font, _fmt_size(style.title_size), style.title_color,
        style.body_font, _fmt_size(style.body_size), style.body_color,
    )
    if style.section_title_size is not None:
        logger.debug("   section title size override: %s", _fmt_size(style.section_title_size))
    if style.content_title_size is not None:
        logger.debug("   content title size override: %s", _fmt_size(style.content_title_size))

    return style
