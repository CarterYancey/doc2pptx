"""Top-level pipeline orchestrator: read → (rewrite) → parse → render → save."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from .llm import (
    DEFAULT_OLLAMA_HOST,
    DEFAULT_OLLAMA_MODEL,
    rewrite_for_pptx_chunked,
)
from .logging_config import logger
from .parser import parse_text_to_deck
from .readers import MAX_TABLE_ROWS_PER_SLIDE, read_document, read_xlsx
from .renderer import (
    _add_content_slide_default,
    _add_section_slide_default,
    _add_slide_from_template,
    _add_table_slide_default,
    _add_table_slide_from_template,
    _add_title_slide_default,
)
from .template import analyze_template


def generate_pptx(
    input_path: str,
    output_path: str,
    template_path: str | None = None,
    title: str | None = None,
    max_bullets: int = 7,
    max_table_rows: int = MAX_TABLE_ROWS_PER_SLIDE,
    use_llm: bool = True,
    ollama_host: str = DEFAULT_OLLAMA_HOST,
    ollama_model: str = DEFAULT_OLLAMA_MODEL,
    llm_prompt_file: str | None = None,
    max_chunk_chars: int = 20000,
    chunk_overlap_sentences: int = 2,
):
    """
    Main entry point: read a document file and produce a .pptx presentation.

    Args:
        input_path:    Path to the input document (.txt, .md, .docx, .pdf, .html, .xlsx, .csv)
        output_path:   Path for the generated .pptx file
        template_path: Optional path to a .pptx template for styling
        title:         Optional override for the presentation title
        max_bullets:   Maximum bullet points per slide before auto-splitting
        max_table_rows: Maximum data rows per table slide (xlsx/csv only)
        use_llm:       If True, rewrite extracted text with a local Ollama server
                       before parsing. Skipped automatically for spreadsheets and
                       when the server is unreachable.
        ollama_host:   Base URL of the Ollama server.
        ollama_model:  Model name to request from Ollama.
        llm_prompt_file: Optional path to a text file containing a custom system
                       prompt to override the default rewrite instructions.
        max_chunk_chars: Maximum characters per LLM chunk. Large inputs are
                       always split on headings/paragraphs so small local models
                       don't overflow their context window.
        chunk_overlap_sentences: Number of trailing sentences from the previous
                       chunk to prepend to the next as orientation context.
    """
    ext = Path(input_path).suffix.lower()
    is_spreadsheet = ext in (".xlsx", ".xls", ".xlsm", ".csv", ".tsv")

    # 1. Read & parse
    logger.info("📄 Reading %s…", input_path)

    if is_spreadsheet:
        deck = read_xlsx(input_path, deck_title=title or "", max_rows_per_slide=max_table_rows)
    else:
        raw_text = read_document(input_path)
        if not raw_text.strip():
            raise ValueError(f"No text content could be extracted from '{input_path}'.")

        if use_llm:
            logger.info("🧠 Rewriting with Ollama (%s) at %s…", ollama_model, ollama_host)
            try:
                system_prompt = None
                if llm_prompt_file:
                    system_prompt = Path(llm_prompt_file).read_text(encoding="utf-8")
                    logger.info("   Using custom system prompt from %s (%d chars)",
                                llm_prompt_file, len(system_prompt))
                rewritten = rewrite_for_pptx_chunked(
                    raw_text,
                    host=ollama_host,
                    model=ollama_model,
                    system_prompt=system_prompt,
                    max_chunk_chars=max_chunk_chars,
                    overlap_sentences=chunk_overlap_sentences,
                )
                if rewritten:
                    raw_text = rewritten
                    logger.info("   LLM rewrite applied (%d chars).", len(rewritten))
                    logger.debug("=== FINAL STITCHED MARKDOWN (%d chars) ===\n%s\n=== END STITCHED MARKDOWN ===",
                                 len(rewritten), rewritten)
                else:
                    logger.warning("   ⚠️  LLM returned empty output; using original text.")
            except Exception as exc:
                logger.warning("   ⚠️  LLM rewrite skipped (%s); using original text.", exc)

        logger.info("🔍 Parsing document structure…")
        deck = parse_text_to_deck(raw_text, deck_title=title or "", max_bullets=max_bullets)

    if not deck.slides:
        raise ValueError("Could not extract any meaningful content for slides.")

    table_count = sum(1 for s in deck.slides if s.slide_type == "table")
    other_count = len(deck.slides) - table_count
    parts = []
    if other_count:
        parts.append(f"{other_count} content")
    if table_count:
        parts.append(f"{table_count} table")
    logger.info("   Found %d slides (%s)", len(deck.slides), ", ".join(parts))
    for idx, sc in enumerate(deck.slides, start=1):
        logger.debug("   slide %d: type=%s level=%d title=%r bullets=%d",
                     idx, sc.slide_type, sc.level, sc.title, len(sc.body_lines))

    # 2. Build the presentation
    if template_path:
        logger.info("🎨 Loading template from %s…", template_path)
        style = analyze_template(template_path)
        prs = Presentation(template_path)

        # Remove any existing slides from the template
        for _ in range(len(prs.slides)):
            sldId = prs.slides._sldIdLst[0]
            rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId is None:
                for attr_key, attr_val in sldId.attrib.items():
                    if attr_key.endswith('}id') or attr_key == 'r:id':
                        rId = attr_val
                        break
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except KeyError:
                    pass
            prs.slides._sldIdLst.remove(sldId)

        logger.info("   Generating slides with template styling…")
        for sc in deck.slides:
            if sc.slide_type == "table":
                _add_table_slide_from_template(prs, style, sc)
            else:
                _add_slide_from_template(prs, style, sc, subtitle=deck.subtitle)
    else:
        logger.info("🎨 Using default styling…")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        for sc in deck.slides:
            if sc.slide_type == "title":
                _add_title_slide_default(prs, sc, subtitle=deck.subtitle)
            elif sc.slide_type == "section":
                _add_section_slide_default(prs, sc)
            elif sc.slide_type == "table":
                _add_table_slide_default(prs, sc)
            else:
                _add_content_slide_default(prs, sc)

    # 3. Save
    prs.save(output_path)
    logger.info("✅ Saved to %s (%d slides)", output_path, len(deck.slides))
    return output_path
