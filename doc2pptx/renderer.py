"""Slide generators — build slides using default styles or a template."""

from __future__ import annotations

import logging

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .logging_config import logger
from .models import SlideContent, TemplateStyle
from .template import _fmt_size, _ph_type_name


# ─── Default (no-template) style constants ───────────────────

DEFAULT_COLORS = {
    "bg_dark": RGBColor(0x1E, 0x27, 0x61),         # navy
    "bg_light": RGBColor(0xFF, 0xFF, 0xFF),        # white
    "title_on_dark": RGBColor(0xFF, 0xFF, 0xFF),   # white text
    "title_on_light": RGBColor(0x1E, 0x27, 0x61),  # navy text
    "body": RGBColor(0x33, 0x33, 0x33),            # dark gray
    "accent": RGBColor(0xCA, 0xDC, 0xFC),          # ice blue
    "subtle": RGBColor(0x64, 0x74, 0x8B),          # muted gray
}
DEFAULT_TITLE_FONT = "Georgia"
DEFAULT_BODY_FONT = "Calibri"


def _apply_font(run, font_name: str, size: Pt, color: RGBColor | None = None, bold: bool = False):
    """Apply font properties to a text run."""
    run.font.name = font_name
    run.font.size = size
    if color:
        run.font.color.rgb = color
    run.font.bold = bold


# ─── Default-style slide renderers ───────────────────────────

def _add_title_slide_default(prs: Presentation, slide_content: SlideContent, subtitle: str = ""):
    """Create a title slide with the default built-in style."""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_dark"]

    # Title
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_content.title
        _apply_font(run, DEFAULT_TITLE_FONT, Pt(40), DEFAULT_COLORS["title_on_dark"], bold=True)

    # Subtitle
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = subtitle or ""
            _apply_font(run, DEFAULT_BODY_FONT, Pt(18), DEFAULT_COLORS["accent"])
            break


def _add_section_slide_default(prs: Presentation, slide_content: SlideContent):
    """Create a section divider slide with default style."""
    # Use a blank layout and build manually
    layout_idx = min(6, len(prs.slide_layouts) - 1)
    slide_layout = prs.slide_layouts[layout_idx]
    logger.info(
        "   slide: type=section title=%r -> default section renderer (blank layout idx=%d %r, navy bg)",
        slide_content.title, layout_idx, slide_layout.name,
    )
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_dark"]

    # Accent line
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(2.2), Inches(1.5), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DEFAULT_COLORS["accent"]
    shape.line.fill.background()

    # Section title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_content.title
    _apply_font(run, DEFAULT_TITLE_FONT, Pt(36), DEFAULT_COLORS["title_on_dark"], bold=True)


def _add_content_slide_default(prs: Presentation, slide_content: SlideContent):
    """Create a content slide with default style."""
    slide_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_light"]

    y_cursor = Inches(0.5)

    # Slide title
    if slide_content.title:
        txBox = slide.shapes.add_textbox(Inches(0.7), y_cursor, Inches(8.6), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_content.title
        _apply_font(run, DEFAULT_TITLE_FONT, Pt(28), DEFAULT_COLORS["title_on_light"], bold=True)
        y_cursor = Inches(1.4)

    # Body content
    if slide_content.body_lines:
        txBox = slide.shapes.add_textbox(
            Inches(0.7), y_cursor, Inches(8.6), Inches(4.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(slide_content.body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)

            # Add a bullet indicator
            p.level = 0
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()

            run = p.add_run()
            run.text = f"  •  {line}"
            _apply_font(run, DEFAULT_BODY_FONT, Pt(15), DEFAULT_COLORS["body"])


# ─── Table helpers (shared by default and template renderers) ───

def _compute_col_widths(headers: list[str], rows: list[list[str]], total_width: float) -> list[float]:
    """
    Compute proportional column widths based on max content length per column.
    Returns widths in inches that sum to total_width.
    """
    n_cols = len(headers)
    if n_cols == 0:
        return []

    max_lens = []
    for c in range(n_cols):
        col_max = len(headers[c])
        for row in rows:
            if c < len(row):
                col_max = max(col_max, len(row[c]))
        max_lens.append(max(col_max, 3))  # minimum 3 chars

    total_chars = sum(max_lens)
    # Proportional, with a minimum width per column
    min_width = 0.8  # inches
    widths = []
    for ml in max_lens:
        w = max(min_width, (ml / total_chars) * total_width)
        widths.append(w)

    # Normalize so they sum to total_width
    scale = total_width / sum(widths)
    widths = [w * scale for w in widths]
    return widths


def _build_table_on_slide(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left: float,
    top: float,
    width: float,
    height: float,
    header_font: str = "Calibri",
    body_font: str = "Calibri",
    header_size: Pt = None,
    body_size: Pt = None,
    header_color: RGBColor | None = None,
    body_color: RGBColor | None = None,
    header_bg: RGBColor | None = None,
    stripe_bg: RGBColor | None = None,
):
    """
    Add a formatted table shape to a slide.
    Handles dynamic column widths and alternating row stripes.
    """
    if header_size is None:
        header_size = Pt(11)
    if body_size is None:
        body_size = Pt(10)
    if header_bg is None:
        header_bg = RGBColor(0x1E, 0x27, 0x61)  # navy
    if header_color is None:
        header_color = RGBColor(0xFF, 0xFF, 0xFF)  # white
    if body_color is None:
        body_color = RGBColor(0x33, 0x33, 0x33)
    if stripe_bg is None:
        stripe_bg = RGBColor(0xF0, 0xF4, 0xF8)  # very light blue-gray

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header row

    # Compute proportional column widths
    col_widths = _compute_col_widths(headers, rows, width)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    table = table_shape.table

    # Set column widths
    for c, w in enumerate(col_widths):
        table.columns[c].width = Inches(w)

    # Disable built-in banding (we'll do our own striping)
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is not None:
        tblPr.set("bandRow", "0")
        tblPr.set("bandCol", "0")
        tblPr.set("firstRow", "0")
        tblPr.set("lastRow", "0")

    # Helper to style a cell
    def _style_cell(cell, text, font_name, font_size, font_color, bg_color=None, bold=False):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        # Reduce internal margins for compact tables
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold

        if bg_color:
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn("a:solidFill"), {})
            srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": str(bg_color)})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

    # Header row
    for c, h in enumerate(headers):
        _style_cell(
            table.cell(0, c), h,
            header_font, header_size, header_color,
            bg_color=header_bg, bold=True,
        )

    # Data rows with alternating stripes
    for r, row_data in enumerate(rows):
        bg = stripe_bg if r % 2 == 1 else None
        for c in range(n_cols):
            val = row_data[c] if c < len(row_data) else ""
            _style_cell(
                table.cell(r + 1, c), val,
                body_font, body_size, body_color,
                bg_color=bg,
            )

    return table_shape


def _add_table_slide_default(prs: Presentation, slide_content: SlideContent):
    """Create a table slide with the default built-in style."""
    slide_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DEFAULT_COLORS["bg_light"]

    y_cursor = 0.4

    # Slide title
    if slide_content.title:
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_cursor), Inches(9.0), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_content.title
        _apply_font(run, DEFAULT_TITLE_FONT, Pt(24), DEFAULT_COLORS["title_on_light"], bold=True)
        y_cursor = 1.1

    # Table
    n_rows = len(slide_content.table_rows)
    # Dynamically size: more rows → smaller row height
    available_height = 5.625 - y_cursor - 0.3  # slide height minus margins
    row_height_estimate = min(0.35, available_height / (n_rows + 1))
    table_height = row_height_estimate * (n_rows + 1)

    _build_table_on_slide(
        slide,
        headers=slide_content.table_headers,
        rows=slide_content.table_rows,
        left=0.5,
        top=y_cursor,
        width=9.0,
        height=table_height,
        header_font=DEFAULT_BODY_FONT,
        body_font=DEFAULT_BODY_FONT,
        header_bg=DEFAULT_COLORS["bg_dark"],
        header_color=DEFAULT_COLORS["title_on_dark"],
        body_color=DEFAULT_COLORS["body"],
        stripe_bg=RGBColor(0xEE, 0xF2, 0xF7),
    )


# ─── Template-style slide renderers ──────────────────────────

def _add_table_slide_from_template(prs: Presentation, style: TemplateStyle, slide_content: SlideContent):
    """Create a table slide using the template's layout and style."""
    layouts = prs.slide_layouts
    layout_idx = min(style.content_layout_idx, len(layouts) - 1)
    layout = layouts[layout_idx]
    n_rows = len(slide_content.table_rows)
    n_cols = len(slide_content.table_headers)
    logger.info(
        "   slide: type=table title=%r -> layout idx=%d %r (%d cols x %d rows)",
        slide_content.title, layout_idx, layout.name, n_cols, n_rows,
    )
    slide = prs.slides.add_slide(layout)

    # Set the title placeholder if available
    title_set = False
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = slide_content.title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = style.title_font
                    if style.title_color:
                        run.font.color.rgb = style.title_color
            title_set = True
            logger.debug(
                "     table title placeholder found: idx=0 name=%r (filled)", ph.name,
            )
            break
    if not title_set:
        logger.debug(
            "     no idx=0 title placeholder on table layout; table title not rendered as placeholder"
        )

    available_height = 4.0
    row_height_estimate = min(0.35, available_height / (n_rows + 1))
    table_height = row_height_estimate * (n_rows + 1)

    # Derive header background from template title color or fall back to navy
    header_bg = style.title_color or RGBColor(0x1E, 0x27, 0x61)

    _build_table_on_slide(
        slide,
        headers=slide_content.table_headers,
        rows=slide_content.table_rows,
        left=0.5,
        top=1.5,
        width=9.0,
        height=table_height,
        header_font=style.body_font,
        body_font=style.body_font,
        header_size=Pt(11),
        body_size=style.body_size if style.body_size else Pt(10),
        header_bg=header_bg,
        header_color=RGBColor(0xFF, 0xFF, 0xFF),
        body_color=style.body_color or RGBColor(0x33, 0x33, 0x33),
    )


def _set_placeholder_text(ph, text: str, fallback_font: str,
                          fallback_size: Pt | None, fallback_color: RGBColor | None):
    """Write ``text`` into a placeholder while preserving the layout's
    existing run formatting where possible.

    The layout placeholder already carries the template designer's font,
    size, weight, and color. Calling ``ph.text = ...`` blows all of that
    away. Instead, reuse the first paragraph's first run (which inherits
    layout styling via the run/paragraph/default chain) and only fill in
    ``font.name``/``font.size``/``font.color`` when the layout itself
    didn't provide one.
    """
    logger.debug(
        "     set placeholder text: idx=%d name=%r text=%r",
        ph.placeholder_format.idx, ph.name, text[:60],
    )
    tf = ph.text_frame
    # Keep the first paragraph (inherits layout pPr); drop the rest.
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    p = tf.paragraphs[0]
    # Preserve the first run if present (keeps its rPr); otherwise make one.
    if p.runs:
        run = p.runs[0]
        # Clear any additional runs so we end up with a single run.
        for extra in list(p.runs[1:]):
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
    run.text = text
    # Only fill in properties the layout didn't already supply.
    applied = []
    if not run.font.name and fallback_font:
        run.font.name = fallback_font
        applied.append(f"name={fallback_font!r}")
    if run.font.size is None and fallback_size:
        run.font.size = fallback_size
        applied.append(f"size={_fmt_size(fallback_size)}")
    try:
        has_color = run.font.color and run.font.color.type is not None
    except (AttributeError, TypeError):
        has_color = False
    if not has_color and fallback_color:
        run.font.color.rgb = fallback_color
        applied.append(f"color={fallback_color}")
    if applied:
        logger.debug(
            "     font fallback applied (layout did not supply): %s", ", ".join(applied),
        )


def _fill_body_placeholder(ph, lines: list[str], style: TemplateStyle):
    """Write ``lines`` into a body placeholder, reusing the layout's
    existing bullet/paragraph formatting for the first line so the
    template's indent and bullet glyph stay intact."""
    logger.debug(
        "     fill body placeholder: idx=%d name=%r %d line(s)",
        ph.placeholder_format.idx, ph.name, len(lines),
    )
    tf = ph.text_frame
    # Drop every paragraph past the first so we can rebuild from a clean slate.
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    first_p = tf.paragraphs[0]
    # Clear runs from the first paragraph but keep its pPr (bullet/indent).
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)

    fallback_counts = {"name": 0, "size": 0, "color": 0}
    for i, line in enumerate(lines):
        if i == 0:
            p = first_p
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if not run.font.name and style.body_font:
            run.font.name = style.body_font
            fallback_counts["name"] += 1
        if run.font.size is None and style.body_size:
            run.font.size = style.body_size
            fallback_counts["size"] += 1
        try:
            has_color = run.font.color and run.font.color.type is not None
        except (AttributeError, TypeError):
            has_color = False
        if not has_color and style.body_color:
            run.font.color.rgb = style.body_color
            fallback_counts["color"] += 1
    if any(fallback_counts.values()):
        logger.debug(
            "     body font fallback applied on %d/%d lines: name=%d size=%d color=%d",
            max(fallback_counts.values()), len(lines),
            fallback_counts["name"], fallback_counts["size"], fallback_counts["color"],
        )


def _add_slide_from_template(prs: Presentation, style: TemplateStyle, slide_content: SlideContent, subtitle: str = ""):
    """Create a slide using the template's layouts and styles."""
    # Section slides without a real Section Header layout would be
    # indistinguishable from content slides; fall back to the default
    # renderer so they still look like section dividers.
    if slide_content.slide_type == "section" and not style.section_layout_found:
        logger.debug(
            "   section slide %r -> default renderer (no Section Header layout in template)",
            slide_content.title,
        )
        _add_section_slide_default(prs, slide_content)
        return

    layouts = prs.slide_layouts

    if slide_content.slide_type == "title":
        layout_idx = style.title_layout_idx
        title_size = style.title_size
    elif slide_content.slide_type == "section":
        layout_idx = style.section_layout_idx
        title_size = style.section_title_size or style.title_size
    else:
        layout_idx = style.content_layout_idx
        title_size = style.content_title_size or style.title_size

    layout_idx = min(layout_idx, len(layouts) - 1)
    layout = layouts[layout_idx]
    logger.info(
        "   slide: type=%s title=%r -> layout idx=%d %r (title_size=%s)",
        slide_content.slide_type, slide_content.title, layout_idx, layout.name, _fmt_size(title_size),
    )
    slide = prs.slides.add_slide(layout)

    # Track what we did with each placeholder so "left-untouched" cases
    # (i.e. stale empty layout boxes) show up in the log.
    ph_actions: list[str] = []
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title placeholder
            _set_placeholder_text(
                ph, slide_content.title,
                fallback_font=style.title_font,
                fallback_size=title_size,
                fallback_color=style.title_color,
            )
            ph_actions.append(
                f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=filled-title"
            )
        elif idx == 1:  # Body/subtitle placeholder
            if slide_content.slide_type == "title" and subtitle:
                _set_placeholder_text(
                    ph, subtitle,
                    fallback_font=style.body_font,
                    fallback_size=style.body_size,
                    fallback_color=style.body_color,
                )
                ph_actions.append(
                    f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=filled-subtitle"
                )
            elif slide_content.body_lines:
                _fill_body_placeholder(ph, slide_content.body_lines, style)
                ph_actions.append(
                    f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=filled-body"
                )
            else:
                ph_actions.append(
                    f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=LEFT-UNTOUCHED (no body/subtitle content)"
                )
        else:
            ph_actions.append(
                f"idx={idx} type={_ph_type_name(ph)} name={ph.name!r} action=LEFT-UNTOUCHED (no handler for this idx)"
            )

    if logger.isEnabledFor(logging.DEBUG):
        if ph_actions:
            logger.debug("     placeholder actions:")
            for line in ph_actions:
                logger.debug("       %s", line)
        else:
            logger.debug("     placeholder actions: (layout has no placeholders)")

    # Fallbacks for minimal layouts that have no title/body placeholders.
    # Without these, H3 slide titles vanish silently because there's nowhere
    # for the title to land.
    title_ph_found = any(ph.placeholder_format.idx == 0 for ph in slide.placeholders)
    body_ph_found = any(ph.placeholder_format.idx == 1 for ph in slide.placeholders)

    body_top = Inches(1.8)
    if not title_ph_found and slide_content.title and slide_content.slide_type != "section":
        logger.debug(
            "     fallback: layout has no title placeholder; adding custom title textbox"
        )
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(1.0))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        trun = tp.add_run()
        trun.text = slide_content.title
        trun.font.name = style.title_font
        trun.font.size = title_size or Pt(28)
        trun.font.bold = True
        if style.title_color:
            trun.font.color.rgb = style.title_color
        body_top = Inches(1.6)

    if not body_ph_found and slide_content.body_lines and slide_content.slide_type == "content":
        logger.debug(
            "     fallback: layout has no body placeholder; adding custom body textbox (%d line(s))",
            len(slide_content.body_lines),
        )
        txBox = slide.shapes.add_textbox(Inches(0.7), body_top, Inches(8.6), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(slide_content.body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"•  {line}"
            run.font.name = style.body_font
            run.font.size = style.body_size
            if style.body_color:
                run.font.color.rgb = style.body_color
            p.space_after = Pt(6)
